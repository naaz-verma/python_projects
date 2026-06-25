"""
Generator for WorldWithWeb_AI_Agents_Course_Deck.pptx
Audience: Working professionals, business owners, advanced students
Run: python sessions/create_ppt_ai_agents_course.py
Output: /Users/macsolutions/Projects/worldwithweb_assets/decks/session_decks/
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x0F, 0x1A)   # deeper navy (more premium than motivation deck)
SURFACE = RGBColor(0x16, 0x18, 0x2E)
CARD    = RGBColor(0x1E, 0x21, 0x3A)
CYAN    = RGBColor(0x00, 0xB4, 0xD8)
GREEN   = RGBColor(0x00, 0xE6, 0x76)
PURPLE  = RGBColor(0xBB, 0x86, 0xFC)
ORANGE  = RGBColor(0xFF, 0x9F, 0x1C)
RED     = RGBColor(0xFF, 0x45, 0x6E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x88, 0x88, 0xAA)
YELLOW  = RGBColor(0xFF, 0xD6, 0x00)

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def txt(slide, text, l, t, w, h, color=WHITE, size=18, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def card(slide, text, l, t, w, h, bg_c=CARD, fg=WHITE, size=16, bold=False):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = bg_c
    sh.line.fill.background()
    tf = sh.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg
    r.font.name = "Calibri"
    return sh


def divider(slide, top=Inches(1.1), color=CYAN):
    sh = slide.shapes.add_shape(1, Inches(0.5), top, Inches(12.33), Inches(0.04))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()


def title_block(slide, heading, sub=""):
    txt(slide, heading,
        Inches(0.5), Inches(0.22), Inches(12.33), Inches(0.9),
        color=CYAN, size=34, bold=True)
    if sub:
        txt(slide, sub,
            Inches(0.5), Inches(1.05), Inches(12.33), Inches(0.45),
            color=GRAY, size=15, italic=True)
    divider(slide, top=Inches(1.05) if not sub else Inches(1.55))


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def flow_box(slide, label, x, y, w=Inches(3.6), h=Inches(0.65), color=CYAN):
    card(slide, label, x, y, w, h, bg_c=CARD, fg=color, size=14, bold=True)


def flow_arrow(slide, x, y):
    txt(slide, "↓", x, y, Inches(0.5), Inches(0.28), color=GRAY, size=13, align=PP_ALIGN.CENTER)


# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def s01_hook(prs):
    """Hook — What if your business ran itself?"""
    slide = _blank(prs)
    bg(slide)

    # Faint ambient text
    for wm, (l, t) in zip(["n8n", "OpenAI", "WhatsApp", "Telegram", "Gemini", "Claude"],
                           [(Inches(0.2), Inches(0.3)), (Inches(10.5), Inches(0.2)),
                            (Inches(0.3), Inches(6.6)), (Inches(10.8), Inches(6.4)),
                            (Inches(0.2), Inches(3.6)), (Inches(11.0), Inches(3.4))]):
        txt(slide, wm, l, t, Inches(2.0), Inches(0.4),
            color=RGBColor(0x22, 0x26, 0x45), size=16, bold=True)

    txt(slide,
        "What If Your Business Ran Itself\nAt 2am — Without You?",
        Inches(0.8), Inches(1.4), Inches(11.73), Inches(2.2),
        color=CYAN, size=42, bold=True, align=PP_ALIGN.CENTER)

    txt(slide,
        "That's not the future.  That's AI Agents running today.",
        Inches(0.8), Inches(3.75), Inches(11.73), Inches(0.75),
        color=WHITE, size=26, align=PP_ALIGN.CENTER)

    # 3 stat chips
    for i, (stat, label, col) in enumerate([
        ("97M+", "New AI jobs by 2025\n(World Economic Forum)", CYAN),
        ("$500B+", "AI automation market\nby 2030", ORANGE),
        ("24/7", "AI agents work —\nno sick days", GREEN),
    ]):
        x = Inches(1.0) + i * Inches(3.9)
        card(slide, stat, x, Inches(4.7), Inches(3.3), Inches(0.7),
             bg_c=SURFACE, fg=col, size=24, bold=True)
        card(slide, label, x, Inches(5.4), Inches(3.3), Inches(0.75),
             bg_c=CARD, fg=WHITE, size=13)

    txt(slide, "WorldWithWeb  ·  Ludhiana's First AI Automation & Agentic AI Course",
        Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.45),
        color=ORANGE, size=13, bold=True, align=PP_ALIGN.CENTER)

    notes(slide, """\
SLIDE 1 — Hook (3 min)

Open with a pause. Let the headline sink in.

"Think about the last time a student enquired at 2am. Did someone reply?
Now imagine — a reply goes out in 3 seconds. Their data is in your CRM.
A WhatsApp lands on their phone. And you didn't lift a finger."

"That's not a dream. That's what we built. And in this course, I'll show you
exactly how — and then teach you to build the same thing for any business."

"This is the AI Agents & Automation course — Ludhiana's first. Let's get into it."
""")
    return slide


def s02_chatbot_vs_agent(prs):
    """Chatbot vs AI Agent — the key distinction"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "Chatbot vs AI Agent — Know the Difference")

    # Left column: Chatbot
    card(slide, "CHATBOT", Inches(0.5), Inches(1.75), Inches(5.6), Inches(0.6),
         bg_c=SURFACE, fg=GRAY, size=18, bold=True)
    for i, line in enumerate([
        "Answers questions only",
        "One conversation at a time",
        "No memory between sessions",
        "Cannot take actions",
        "Cannot connect to other tools",
    ]):
        card(slide, f"✗  {line}", Inches(0.5), Inches(2.45) + i * Inches(0.72),
             Inches(5.6), Inches(0.6), bg_c=CARD, fg=RED, size=15)

    # Right column: Agent
    card(slide, "AI AGENT", Inches(7.2), Inches(1.75), Inches(5.6), Inches(0.6),
         bg_c=SURFACE, fg=CYAN, size=18, bold=True)
    for i, line in enumerate([
        "Understands & decides",
        "Runs autonomously 24/7",
        "Remembers context & history",
        "Sends emails, WhatsApp, updates sheets",
        "Connects to ANY tool via API",
    ]):
        card(slide, f"✓  {line}", Inches(7.2), Inches(2.45) + i * Inches(0.72),
             Inches(5.6), Inches(0.6), bg_c=CARD, fg=GREEN, size=15)

    txt(slide, "VS", Inches(6.1), Inches(3.6), Inches(1.1), Inches(0.7),
        color=ORANGE, size=28, bold=True, align=PP_ALIGN.CENTER)

    card(slide, "ChatGPT is a chatbot. An AI Agent ACTS — it does things, connects tools, and completes tasks end-to-end.",
         Inches(0.5), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=YELLOW, size=14)

    notes(slide, """\
SLIDE 2 — Chatbot vs Agent (3 min)

"Everyone here has used ChatGPT. So let's be clear about what the difference actually is."

"A chatbot answers your question. Done. It waits for the next question."
[Point to left column]

"An AI agent? It understands your goal, makes decisions, and ACTS.
It can send emails. Update spreadsheets. Post to Instagram. Message on WhatsApp.
All without you being in the loop."
[Point to right column]

"When a student enquiries on your website at 2am — a chatbot says 'Thanks, we'll call you.'
An AI agent replies, saves the data, sends WhatsApp, logs it in CRM, and emails your team."

"THAT is the difference we're teaching in this course."
""")
    return slide


def s03_what_agents_do(prs):
    """What agentic AI can actually do today"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "What AI Agents Can Do Right Now", "Not in 5 years. Today.")

    capabilities = [
        ("📧  Email", "Read, classify, respond, forward — based on rules you set", CYAN),
        ("💬  WhatsApp", "Send messages, receive replies, qualify leads automatically", GREEN),
        ("📊  Spreadsheets", "Read data, update rows, generate reports, trigger actions", CYAN),
        ("📱  Social Media", "Write captions, generate images, post on schedule", PURPLE),
        ("🔔  CRM / Sheets", "Log every lead, update stages, alert your team", ORANGE),
        ("🌐  Web Scraping", "Monitor competitors, prices, job boards, news — automatically", GREEN),
        ("📞  Voice AI", "Answer calls, qualify leads, book appointments via voice", PURPLE),
        ("📄  Documents", "Extract data from PDFs, invoices, forms — zero manual entry", ORANGE),
    ]

    cw, ch = Inches(5.8), Inches(0.62)
    gap = Inches(0.18)
    sx = Inches(0.5)
    sy = Inches(1.75)

    for i, (icon_title, desc, color) in enumerate(capabilities):
        row, col = i // 2, i % 2
        x = sx + col * (cw + Inches(0.73))
        y = sy + row * (ch + gap)
        card(slide, icon_title, x, y, Inches(2.0), ch, bg_c=SURFACE, fg=color, size=14, bold=True)
        card(slide, desc, x + Inches(2.05), y, Inches(3.75), ch, bg_c=CARD, fg=WHITE, size=13)

    notes(slide, """\
SLIDE 3 — What Agents Do (3 min)

Walk through each row quickly. The key is breadth — students should feel like
"wait, AI can do ALL of that?"

"Email: you set the rules — if someone emails about Python courses, the agent
categorises it, replies with course details, and logs it."

"WhatsApp: not just sending. Receiving replies and understanding them."

"Social media: you'll see this in Demo #3 — fully automated content pipeline."

"Documents: one of the highest-paying freelance skills right now.
Businesses have thousands of PDFs they need data extracted from."

"This is what the course covers. All of this becomes something you can build and sell."
""")
    return slide


def s04_market_opportunity(prs):
    """The market opportunity"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "The Market Is Paying — Right Now")

    stats = [
        ("$500B+", "Global automation\nmarket by 2030", CYAN),
        ("67%", "SMBs want AI automation\nbut lack the talent", ORANGE),
        ("₹50K–₹2L", "Per project rate for\nAI automation builds", GREEN),
        ("₹30K–₹80K", "Monthly retainer for\nagent maintenance", PURPLE),
    ]
    sw = Inches(2.85)
    for i, (val, label, col) in enumerate(stats):
        x = Inches(0.5) + i * (sw + Inches(0.35))
        card(slide, val, x, Inches(1.75), sw, Inches(1.0), bg_c=SURFACE, fg=col, size=26, bold=True)
        card(slide, label, x, Inches(2.75), sw, Inches(0.85), bg_c=CARD, fg=WHITE, size=14)

    industries = [
        ("Real Estate", "Lead qualification agents, property alert bots"),
        ("Healthcare", "Appointment booking, patient follow-up, report extraction"),
        ("Education", "Enquiry handling, follow-up sequences, result notifications"),
        ("E-commerce", "Order updates, review requests, inventory alerts"),
        ("Marketing", "Content pipeline, lead nurturing, analytics reports"),
        ("Finance", "Invoice processing, payment reminders, fraud alerts"),
    ]
    txt(slide, "Industries hiring automation consultants today:",
        Inches(0.5), Inches(3.8), Inches(12.33), Inches(0.4), color=GRAY, size=14)

    iw, ih = Inches(3.85), Inches(0.65)
    for i, (title, desc) in enumerate(industries):
        row, col = i // 2, i % 2
        x = Inches(0.5) + col * (iw + Inches(0.6))
        y = Inches(4.3) + row * (ih + Inches(0.15))
        card(slide, title, x, y, Inches(1.3), ih, bg_c=SURFACE, fg=ORANGE, size=13, bold=True)
        card(slide, desc, x + Inches(1.35), y, Inches(2.5), ih, bg_c=CARD, fg=WHITE, size=12)

    notes(slide, """\
SLIDE 4 — Market (3 min)

"Let me show you the size of the opportunity."

"$500 billion — that's the automation market by 2030. That's not AI in general — that's
specifically workflow and process automation."

"67% of small businesses WANT to automate but don't have the technical talent.
That gap is where you come in."

"₹50,000 to ₹2 lakhs per project. Real numbers. I'll show you real client rates later."

"And the industries are everywhere — real estate, clinics, schools, e-commerce.
Every single one of these businesses needs exactly what we'll teach you to build."
""")
    return slide


def s05_poc1(prs):
    """POC #1 — Contact Form AI Agent (live WorldWithWeb system)"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "POC #1 — Contact Form AI Agent", "Live system running for WorldWithWeb right now")

    steps = [
        ("Student fills form on website",   CYAN),
        ("Webhook triggers instantly",       GRAY),
        ("Auto-generates Enquiry ID + date", PURPLE),
        ("Saves to Google Sheets (CRM)",     ORANGE),
        ("Sends welcome email to student",   CYAN),
        ("Alerts internal team via email",   ORANGE),
        ("Sends WhatsApp to student",        GREEN),
    ]
    bw = Inches(5.2)
    bh = Inches(0.58)
    sx = Inches(4.0)
    sy = Inches(1.75)
    gap = Inches(0.67)

    for i, (label, color) in enumerate(steps):
        y = sy + i * gap
        card(slide, label, sx, y, bw, bh, bg_c=CARD, fg=color, size=14, bold=(i in [0, 6]))
        if i < len(steps) - 1:
            txt(slide, "↓", sx + Inches(2.3), y + bh,
                Inches(0.6), Inches(0.15), color=GRAY, size=11, align=PP_ALIGN.CENTER)

    # Left annotation
    card(slide, "⚙️  Tools Used",
         Inches(0.3), Inches(1.75), Inches(3.4), Inches(0.5),
         bg_c=SURFACE, fg=GRAY, size=13, bold=True)
    for i, tool in enumerate(["n8n (workflow engine)", "Google Sheets API",
                               "Gmail API", "Evolution API (WhatsApp)", "JavaScript (Code node)"]):
        card(slide, tool, Inches(0.3), Inches(2.35) + i * Inches(0.65),
             Inches(3.4), Inches(0.55), bg_c=CARD, fg=WHITE, size=13)

    card(slide, "Response time: < 3 seconds  ·  Runs 24/7  ·  Zero manual work per lead",
         Inches(0.3), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=GREEN, size=14)

    notes(slide, """\
SLIDE 5 — POC #1 (4 min)

"This is not a hypothetical. This is the actual system running for WorldWithWeb right now."

"Every time a student fills the contact form on our website:
- In under 3 seconds, they get a professional email with their enquiry ID.
- Simultaneously, a WhatsApp message lands on their phone.
- Our team gets an internal email alert.
- The data is already in our Google Sheet, timestamped, with a unique ID.
- Nothing was done manually."

"The workflow is built in n8n. No code except one small JavaScript snippet
for generating the Enquiry ID. Everything else is drag and drop."

"This is a real POC that you'll understand how to build in Module 2 of this course."

[If showing on laptop: Open n8n and show the actual workflow live]
""")
    return slide


def s06_poc2(prs):
    """POC #2 — Smart Follow-up Sequence"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "POC #2 — Smart Lead Follow-up Engine", "Automated 30-day nurture sequence — zero manual effort")

    # Timeline spine
    stages = [
        ("Day 0",  "Enquiry received → CRM logged",               CYAN),
        ("Day 3",  "Email: 'Do you have questions?'",              PURPLE),
        ("Day 7",  "Email: 'New batch starting soon'",             ORANGE),
        ("Day 14", "Email: 'Only a few seats left'",               RED),
        ("Day 30", "Email: 'Early bird offer — last chance'",      GREEN),
        ("Anytime","Status → Converted / Closed → Stop sequence", GRAY),
    ]

    bw = Inches(7.5)
    bh = Inches(0.62)
    sx = Inches(2.5)
    sy = Inches(1.75)

    for i, (day, label, color) in enumerate(stages):
        y = sy + i * Inches(0.82)
        card(slide, day, sx, y, Inches(1.1), bh, bg_c=SURFACE, fg=color, size=14, bold=True)
        card(slide, label, sx + Inches(1.15), y, Inches(6.35), bh, bg_c=CARD, fg=WHITE, size=14)
        if i < len(stages) - 1:
            txt(slide, "↓", sx + Inches(0.3), y + bh,
                Inches(0.5), Inches(0.22), color=GRAY, size=11, align=PP_ALIGN.CENTER)

    # Right annotation
    card(slide, "⚙️  How It Works", Inches(10.3), Inches(1.75), Inches(2.7), Inches(0.5),
         bg_c=SURFACE, fg=GRAY, size=12, bold=True)
    for i, note in enumerate(["Runs daily at 10am", "Reads all enquiries", "Calculates days since",
                               "Smart routing logic", "Updates CRM after send", "Skips converted leads"]):
        card(slide, note, Inches(10.3), Inches(2.35) + i * Inches(0.63),
             Inches(2.7), Inches(0.52), bg_c=CARD, fg=WHITE, size=12)

    card(slide, "One workflow replaces a full-time follow-up team. Personalised emails to every lead. Automatically.",
         Inches(0.3), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=ORANGE, size=14)

    notes(slide, """\
SLIDE 6 — POC #2 (4 min)

"Now here's the follow-up system. This one might be the most valuable for any business."

"A lead comes in. From that moment, the system watches.
Day 3 — they get a friendly email asking if they have questions.
Day 7 — they hear about a new batch starting.
Day 14 — urgency: limited seats.
Day 30 — a special early bird offer."

"Every email is personalised with their name and course interest — pulled from the sheet."

"The system runs every morning at 10am. It calculates exactly where each lead is in
the funnel, sends the right email, then updates the sheet so it doesn't double-send."

"If a lead converts? It stops automatically. If they ask to be removed? One column update."

"This replaced manual follow-up calls and SMS for WorldWithWeb. Zero missed leads."

[Show n8n workflow live if possible — students see the Switch node routing]
""")
    return slide


def s07_poc3(prs):
    """POC #3 — AI Social Media Manager (Week 3 course build)"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "POC #3 — AI Social Media Manager", "Idea → caption + image → Telegram approval → scheduled post")

    steps = [
        ("Content idea in Google Sheets",    "Topic · Platform · Tone · Date",              CYAN,   GRAY),
        ("GPT-4o generates post caption",    "Platform-specific, branded tone",             ORANGE, GRAY),
        ("DALL-E 3 generates post image",    "Auto-prompt derived from caption",            PURPLE, GRAY),
        ("Telegram: preview sent to admin",  "Caption + image — approve with one tap",      YELLOW, GRAY),
        ("On approval: schedule to socials", "Instagram · LinkedIn · Facebook via Buffer",  GREEN,  GRAY),
        ("Google Sheets: status = posted",   "Full audit trail — every post logged",        CYAN,   GRAY),
        ("WhatsApp daily summary to admin",  "What went out today — one message",           GREEN,  GRAY),
    ]

    bw = Inches(4.7)
    bh = Inches(0.6)
    sx = Inches(0.4)
    sy = Inches(1.75)

    for i, (main_l, sub_l, main_c, sub_c) in enumerate(steps):
        y = sy + i * Inches(0.75)
        card(slide, main_l, sx, y, bw, bh, bg_c=CARD, fg=main_c, size=14, bold=True)
        card(slide, sub_l, sx + bw + Inches(0.1), y, Inches(3.1), bh,
             bg_c=SURFACE, fg=WHITE, size=13)
        if i < len(steps) - 1:
            txt(slide, "↓", sx + Inches(2.1), y + bh,
                Inches(0.6), Inches(0.18), color=GRAY, size=11, align=PP_ALIGN.CENTER)

    # Right panel
    card(slide, "⚙️  Stack", Inches(8.1), Inches(1.75), Inches(4.9), Inches(0.5),
         bg_c=SURFACE, fg=GRAY, size=13, bold=True)
    for i, tool in enumerate(["n8n (scheduler + logic)", "Google Sheets (content calendar)",
                               "OpenAI GPT-4o (captions)", "DALL-E 3 (images)",
                               "Telegram Bot API (approval gate)",
                               "Buffer API (Instagram/LinkedIn/Facebook)",
                               "Evolution API (WhatsApp report)"]):
        card(slide, tool, Inches(8.1), Inches(2.35) + i * Inches(0.62),
             Inches(4.9), Inches(0.52), bg_c=CARD, fg=WHITE, size=13)

    card(slide, "Human stays in control — Telegram approval before anything goes live. You plan, AI executes, you approve.",
         Inches(0.3), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=PURPLE, size=14)

    notes(slide, """\
SLIDE 7 — POC #3: AI Social Media Manager (4 min)

"This is the Week 3 build in the marketing track — and the most requested demo."

"Here's the key difference from a basic auto-poster: there's a HUMAN approval gate."

"The AI generates the caption and image. Before anything goes live,
you get a Telegram message with a preview — caption + image — right on your phone.
One tap: approve. The post schedules. Reject: it's discarded. You stay in control."

"Walk through the flow:
You write a one-line content idea in a Google Sheet.
n8n picks it up, sends it to GPT-4o which writes a platform-specific caption.
DALL-E 3 generates a matching image.
A Telegram message lands on your phone with the full preview.
You approve — it schedules to Instagram, LinkedIn, and Facebook simultaneously.
Sheet updates, WhatsApp summary goes out."

"The Telegram approval step is what makes this professional-grade.
Clients trust automations that keep a human in the loop."

"You'll build this from scratch in Week 3 of the course."
""")
    return slide


def s08_usecase_business(prs):
    """Use cases — Business owners"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "What Businesses Are Paying For Today")

    usecases = [
        ("🏡  Real Estate",
         ["Auto-reply to property enquiries", "Send property PDFs on WhatsApp",
          "Schedule site visit reminders", "Notify team of hot leads"], CYAN),
        ("🏥  Clinics & Healthcare",
         ["Appointment booking via WhatsApp", "Patient follow-up reminders",
          "Extract data from prescription PDFs", "No-show re-booking agent"], GREEN),
        ("🛒  E-commerce",
         ["Order status WhatsApp updates", "Review request automation",
          "Abandoned cart follow-up", "Low stock / reorder alerts"], ORANGE),
        ("🎓  Coaching / Training",
         ["Enquiry → instant reply + CRM", "Batch reminders to enrolled students",
          "Certificate generation & delivery", "30-day lead nurture sequence"], PURPLE),
    ]

    uw = Inches(2.9)
    for i, (title, points, color) in enumerate(usecases):
        row, col = i // 2, i % 2
        x = Inches(0.4) + col * (uw + Inches(3.75))
        y = Inches(1.75) + row * Inches(2.3)
        card(slide, title, x, y, Inches(6.43), Inches(0.55),
             bg_c=SURFACE, fg=color, size=16, bold=True)
        for j, pt in enumerate(points):
            card(slide, f"→  {pt}", x, y + Inches(0.6) + j * Inches(0.42),
                 Inches(6.43), Inches(0.38), bg_c=CARD, fg=WHITE, size=13)

    card(slide, "Every one of these is a project you can freelance. Businesses pay ₹30,000–₹1,50,000 per automation.",
         Inches(0.4), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=YELLOW, size=14)

    notes(slide, """\
SLIDE 8 — Business Use Cases (3 min)

"Let me show you who's paying for this skill right now."

Walk through each industry quickly:

"Real estate — every property agent gets 50 enquiries a week.
They can't reply to all of them. An AI agent that replies instantly and sends the property PDF?
Worth ₹50,000 to them easily."

"Healthcare — appointment booking alone saves a receptionist 3 hours a day.
Clinics are paying for this."

"E-commerce — order updates, review requests, abandoned cart recovery.
These are proven revenue increases. Shops pay monthly retainers for this."

"Education — you've just seen our own system. WorldWithWeb uses exactly this.
Any coaching centre would pay for this."

"One agent per business. ₹30,000 to ₹1,50,000 per build. That's the market."
""")
    return slide


def s09_usecase_marketer(prs):
    """Use cases — Digital marketers & freelancers"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "AI Automation for Marketers & Freelancers")

    categories = [
        ("📊  Lead Generation", CYAN, [
            "Scrape LinkedIn leads → verify → add to CRM",
            "Auto-DM new followers with offer",
            "Webinar registrant → WhatsApp follow-up",
            "Form fill → instant qualification call booking",
        ]),
        ("✍️  Content Creation", ORANGE, [
            "Blog idea → draft → WordPress publish → social post",
            "YouTube transcript → LinkedIn article + Twitter thread",
            "Competitor analysis → content gap report",
            "Monthly content calendar — generated in 30 seconds",
        ]),
        ("📈  Reporting & Analytics", PURPLE, [
            "Google Analytics → weekly summary email",
            "Ad spend vs leads → automated ROI report",
            "Social engagement → daily dashboard update",
            "Client report PDF — auto-generated monthly",
        ]),
    ]

    cw = Inches(3.9)
    for i, (title, color, points) in enumerate(categories):
        x = Inches(0.4) + i * (cw + Inches(0.31))
        card(slide, title, x, Inches(1.75), cw, Inches(0.55),
             bg_c=SURFACE, fg=color, size=15, bold=True)
        for j, pt in enumerate(points):
            card(slide, pt, x, Inches(2.4) + j * Inches(1.0),
                 cw, Inches(0.85), bg_c=CARD, fg=WHITE, size=13)

    card(slide, "Freelance rate for AI automation: ₹500–₹2,000/hour. Retainers: ₹25,000–₹80,000/month.",
         Inches(0.4), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=GREEN, size=14)

    notes(slide, """\
SLIDE 9 — Marketer Use Cases (2 min)

"If you're a digital marketer — or want to become one — here's what AI automation adds to your toolkit."

"Lead generation agents are massive. Scraping, enriching, and automating outreach
is something every marketing agency wants but few can build."

"Content automation: one of the highest-ROI uses of AI right now.
The YouTube → LinkedIn → Twitter pipeline alone saves 4 hours per piece of content."

"Reporting: clients love dashboards. Agencies charge ₹10,000–₹20,000/month
just to maintain automated reporting systems."

"Add these to your freelance service offerings and you double your rates immediately."
""")
    return slide


def s10_what_youll_build(prs):
    """4 sector builds + capstone"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "What You'll Build — One Project Per Sector", "4 live automations + your own capstone  ·  Live, not slides")

    projects = [
        {
            "track": "Business & Sales",
            "title": "AI Lead Engine",
            "desc": "AI agent qualifies & replies on WhatsApp 24/7, follows up automatically, sends owner a daily digest",
            "tools": "n8n · WhatsApp · OpenAI · Google Sheets · Gmail",
            "value": "₹30K–₹60K per client build",
            "color": CYAN,
        },
        {
            "track": "Digital Marketing & Social",
            "title": "AI Social Media Manager",
            "desc": "Idea → AI caption + image → Telegram approval → scheduled to Instagram / LinkedIn / Facebook",
            "tools": "n8n · OpenAI · DALL-E · Telegram · Buffer",
            "value": "₹40K–₹80K per project",
            "color": PURPLE,
        },
        {
            "track": "Data Science & ML",
            "title": "Deploy-Your-Model + Data Pipeline",
            "desc": "Wrap a trained ML model behind a webhook → live predictions on WhatsApp + scheduled data pipeline",
            "tools": "n8n · Webhooks · OpenAI · Google Sheets · WhatsApp",
            "value": "₹50K–₹1.2L per build",
            "color": ORANGE,
        },
        {
            "track": "Personal & Career",
            "title": "Your AI Assistant",
            "desc": "Aggregates & filters job listings to your fit, tailors resume & cover letter per role, sends reminders",
            "tools": "n8n · OpenAI · Gmail · Google Sheets · Calendar",
            "value": "Portfolio project + personal productivity",
            "color": GREEN,
        },
    ]

    ph_track = Inches(0.42)
    ph_title = Inches(0.48)
    ph_desc  = Inches(0.52)
    ph_tools = Inches(0.38)
    block_h  = ph_track + ph_title + ph_desc + ph_tools + Inches(0.12)
    sy = Inches(1.75)

    for i, p in enumerate(projects):
        y = sy + i * (block_h + Inches(0.1))
        card(slide, p["track"], Inches(0.4), y, Inches(2.4), ph_track,
             bg_c=p["color"], fg=BG, size=12, bold=True)
        card(slide, p["title"], Inches(2.85), y, Inches(9.85), ph_track,
             bg_c=SURFACE, fg=p["color"], size=14, bold=True)
        card(slide, p["desc"], Inches(0.4), y + ph_track, Inches(8.5), ph_desc + ph_title,
             bg_c=CARD, fg=WHITE, size=13)
        card(slide, p["tools"], Inches(9.0), y + ph_track, Inches(4.1), ph_desc + ph_title,
             bg_c=CARD, fg=GRAY, size=12)
        card(slide, p["value"], Inches(0.4), y + ph_track + ph_desc + ph_title, Inches(12.6), ph_tools,
             bg_c=SURFACE, fg=GREEN, size=12)

    card(slide, "+ Capstone: take the track that fits YOUR goal to full depth — ship it for your business, a client, or your portfolio",
         Inches(0.4), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=YELLOW, size=14, bold=True)

    notes(slide, """\
SLIDE 10 — 4 Sector Builds + Capstone (4 min)

"The course is structured around 4 real sectors. You build one automation in each.
These aren't exercises — they're live, deployable systems."

Track 1 — Business & Sales:
"The AI Lead Engine. An AI agent on WhatsApp that qualifies leads, replies instantly,
follows up across 30 days, and sends you a morning digest. You saw our version live."

Track 2 — Digital Marketing:
"The AI Social Media Manager. The full pipeline — idea to approved, scheduled post.
The Telegram approval gate keeps you in control. Builds to Instagram, LinkedIn, Facebook."

Track 3 — Data Science / ML:
"Deploy-Your-Model. If you've trained any ML model — in college or on Kaggle —
you'll wrap it behind a webhook so it answers live predictions on WhatsApp.
Plus a scheduled data pipeline that fetches, cleans, and AI-summarises data."

Track 4 — Career:
"Your personal AI job assistant. It monitors job boards for roles that match your profile,
tailors your resume and cover letter for each one, and sends you reminders.
Build it for yourself first — then offer it as a service."

"Then the capstone: you pick the track that fits your goal and take it to full depth.
A complete, real automation you ship — for your own business, a client, or your portfolio."
""")
    return slide


def s11_tools(prs):
    """Tools you'll master"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "Tools You'll Master", "The exact stack used in the course — no fluff")

    tools = [
        ("n8n",             "Core workflow engine — the tool\nthe whole course is built on",          CYAN),
        ("OpenAI / Gemini\n/ Claude", "GPT-4o, DALL-E 3, Gemini —\ncall any LLM from your workflows", GREEN),
        ("AI Agent Node",   "n8n's built-in agent — AI that\ndecides and takes actions, not just chats", ORANGE),
        ("WhatsApp API",    "Evolution API — the same setup\nWorldWithWeb runs right now",             CYAN),
        ("Telegram",        "Approval gates, notifications,\nbot triggers inside workflows",           PURPLE),
        ("Google Sheets",   "No-cost CRM + data store —\nread, write, update from n8n",               GREEN),
        ("Gmail",           "Automated emails: confirmations,\nfollow-ups, internal alerts",           ORANGE),
        ("Webhooks & HTTP", "How the internet talks — the\nfoundation of every integration",           PURPLE),
        ("Calendar",        "Schedule triggers, reminders,\ncalendar-driven automations",              CYAN),
    ]

    tw = Inches(3.85)
    th = Inches(1.2)
    gap_x = Inches(0.28)
    gap_y = Inches(0.22)
    sx = Inches(0.4)
    sy = Inches(1.7)
    cols = 3

    for i, (name, desc, color) in enumerate(tools):
        row, col = i // cols, i % cols
        x = sx + col * (tw + gap_x)
        y = sy + row * (th + gap_y)
        card(slide, name, x, y, tw, Inches(0.52), bg_c=SURFACE, fg=color, size=13, bold=True)
        card(slide, desc, x, y + Inches(0.52), tw, Inches(0.62), bg_c=CARD, fg=WHITE, size=12)

    card(slide, "You'll understand WHY each tool exists — not just how to click buttons. That's the WorldWithWeb difference.",
         Inches(0.4), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=YELLOW, size=14)

    notes(slide, """\
SLIDE 11 — Tools (2 min)

"Here's exactly what you'll use in the course."

"n8n is the backbone — everything runs through it.
Open-source, self-hostable, more powerful than Zapier at any scale."

"You don't just use one LLM — you'll connect OpenAI, Google Gemini, AND Anthropic Claude.
You'll know how to choose the right model for the right task."

"The AI Agent node is the crown jewel. This is what separates n8n from other tools.
It's an AI that can USE other tools — search the web, query a sheet, send a message —
all in one step. That's what 'agentic AI' means in practice."

"WhatsApp — the Evolution API that we run ourselves for WorldWithWeb.
Real experience, not a sandbox demo."

"Telegram is your approval and notification channel. It's the human-in-the-loop layer."

"Webhooks and HTTP requests: once you understand these, you can integrate
literally any tool on the internet — regardless of whether n8n has a built-in node for it."
""")
    return slide


def s12_course_structure(prs):
    """Course structure — 4 weeks / 16 sessions"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "Course Structure — 1 Month · 16 Hands-On Sessions")

    modules = [
        ("Week 1",
         "Foundations & First Automation",
         "Automation & agents · n8n setup · Triggers · Google Sheets as mini-CRM · Multi-channel output\nBuild: Intake-to-response automation — the base every track builds on",
         CYAN),
        ("Week 2",
         "Agentic AI + Business & Sales Track",
         "Plug LLMs into n8n · Prompt design · AI Agent node + tools + knowledge base · Lead scoring & multi-stage follow-up\nBuild: AI Lead Engine — qualifies on WhatsApp 24/7 · follows up · daily digest to owner",
         PURPLE),
        ("Week 3",
         "Marketing & Data Science / ML Tracks",
         "Marketing: captions + images + Telegram approval + scheduling  ·  ML: model-behind-webhook + data pipeline\nBuild: AI Social Media Manager  +  Deploy-Your-Model (live WhatsApp predictions + data pipeline)",
         ORANGE),
        ("Week 4",
         "Career Track + Your Capstone",
         "Career: AI assistant filters jobs, tailors resume/cover letter · Packaging automations as a service · Client pitching\nBuild: Your AI Job Assistant  +  Capstone — full automation in your chosen field, shipped",
         GREEN),
    ]

    sy = Inches(1.75)
    mh = Inches(1.12)
    gap = Inches(0.12)

    for i, (week, title, content, color) in enumerate(modules):
        y = sy + i * (mh + gap)
        card(slide, week, Inches(0.4), y, Inches(1.1), mh,
             bg_c=SURFACE, fg=color, size=14, bold=True)
        card(slide, title, Inches(1.55), y, Inches(3.2), mh,
             bg_c=SURFACE, fg=color, size=14, bold=True)
        card(slide, content, Inches(4.8), y, Inches(8.13), mh,
             bg_c=CARD, fg=WHITE, size=12)

    card(slide, "4 sessions/week · 1 month · Practical + project-based · Recorded sessions included",
         Inches(0.4), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=GRAY, size=14)

    notes(slide, """\
SLIDE 12 — Course Structure (3 min)

"1 month. 16 sessions. 4 sessions per week. Intensive and practical."

"Week 1: Foundations. Before you build anything, you need to understand the plumbing —
how data moves between tools, what webhooks are, how expressions work in n8n.
By the end of Week 1 you have a working automation live."

"Week 2: This is where agentic AI comes in. You plug an LLM into n8n.
Then you use the AI Agent node — n8n's most powerful feature.
An agent that can USE tools: look up a spreadsheet, send a WhatsApp, check a website.
Build: the AI Lead Engine. It's live on WhatsApp and runs 24/7."

"Week 3: Two tracks in one week.
Marketing — the Social Media Manager you saw in POC #3.
ML — if you've built any model, you'll deploy it behind a webhook for live predictions."

"Week 4: Career automation + the capstone.
The capstone is your own real project — your business, a client's business, or your portfolio.
You also learn how to price and pitch this work to clients."

"One month. Four working automations. Your capstone shipped."
""")
    return slide


def s13_freelance_income(prs):
    """Freelance income potential"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "Freelance Income Potential", "Conservative estimates based on current market rates in India")

    services = [
        ("AI Lead Engine",           "₹15,000–₹50,000",  "per project",          CYAN),
        ("Social Media Manager",     "₹20,000–₹60,000",  "per project",          PURPLE),
        ("Deploy-Your-Model",        "₹25,000–₹80,000",  "per project",          ORANGE),
        ("Full Business Automation", "₹50,000–₹1,50,000","per project",          GREEN),
        ("Workflow Retainer",        "₹10,000–₹40,000",  "per client/month",     CYAN),
        ("Data Pipeline",            "₹20,000–₹60,000",  "per project",          ORANGE),
    ]

    sw = Inches(3.8)
    sh = Inches(0.95)
    sx = Inches(0.4)
    sy = Inches(1.75)

    for i, (service, rate, unit, color) in enumerate(services):
        row, col = i // 2, i % 2
        x = sx + col * (sw + Inches(0.5) + Inches(2.12))
        y = sy + row * (sh + Inches(0.22))
        card(slide, service, x, y, sw, Inches(0.5), bg_c=SURFACE, fg=WHITE, size=13)
        card(slide, rate, x + sw + Inches(0.1), y, Inches(2.5), Inches(0.5),
             bg_c=CARD, fg=color, size=14, bold=True)
        card(slide, unit, x + sw + Inches(0.1), y + Inches(0.5), Inches(2.5), Inches(0.42),
             bg_c=SURFACE, fg=GRAY, size=12)

    card(slide,
         "Automations sell for ₹15K–₹50K each (course figure). 1 sold project recovers the entire course fee.",
         Inches(0.4), Inches(5.9), Inches(12.33), Inches(0.6),
         bg_c=SURFACE, fg=GREEN, size=15, bold=True)

    card(slide, "Clients pay for outcomes — time saved, leads captured, hours freed. Price the value, not your hours.",
         Inches(0.4), Inches(6.6), Inches(12.33), Inches(0.55),
         bg_c=CARD, fg=GRAY, size=13)

    notes(slide, """\
SLIDE 13 — Freelance Income (3 min)

"Let me be direct about money — because that's why you're here."

"These rates are conservative. They're based on actual projects, not aspirational numbers."

"Lead capture agent: any business gets value from this immediately.
They see ROI in Week 1. ₹30,000–₹60,000 is very reasonable."

"Maintenance is the game-changer. You build once, you maintain monthly.
3 clients paying ₹25,000/month = ₹75,000 recurring. Passive-ish income."

"The compound effect: build a portfolio of 5 automations over 3 months.
Show them in a Notion portfolio. Your rate goes up. Referrals start coming in."

"Important caveat: I'm not promising you'll earn this in Month 1.
I'm showing you what the ceiling looks like and what the path is.
The course gives you the skills. You provide the hustle."
""")
    return slide


def s14_career_paths(prs):
    """Career paths"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "Career Paths After This Course")

    paths = [
        ("AI Automation\nConsultant",     "Freelance / agency. Solve business\nproblems with automation.",  "₹6–20 LPA\nor freelance", CYAN),
        ("AI Solutions\nEngineer",        "In-house at companies building\nand maintaining AI workflows.", "₹8–25 LPA",              PURPLE),
        ("No-Code / AI\nProduct Builder", "Build SaaS tools and automation\nproducts for a niche market.",  "Unlimited\n(founder path)", ORANGE),
        ("Digital Marketing\nAI Lead",    "Run AI-powered marketing for\nbusinesses. Strategy + execution.", "₹5–18 LPA",             GREEN),
    ]

    pw = Inches(2.9)
    ph = Inches(3.2)
    gap = Inches(0.33)
    sx = Inches(0.4)

    for i, (title, desc, salary, color) in enumerate(paths):
        x = sx + i * (pw + gap)
        card(slide, title, x, Inches(1.75), pw, Inches(0.85),
             bg_c=SURFACE, fg=color, size=16, bold=True)
        card(slide, desc, x, Inches(2.65), pw, Inches(1.1),
             bg_c=CARD, fg=WHITE, size=13)
        card(slide, "Income:", x, Inches(3.8), pw, Inches(0.45),
             bg_c=SURFACE, fg=GRAY, size=12)
        card(slide, salary, x, Inches(4.25), pw, Inches(0.7),
             bg_c=SURFACE, fg=color, size=15, bold=True)

    card(slide, "Prerequisite paths: Python + Data Science track → adds AI Engineer and ML Engineer to your options.",
         Inches(0.4), Inches(5.2), Inches(12.33), Inches(0.55),
         bg_c=CARD, fg=GRAY, size=13)

    card(slide, "This course alone — no prior coding required — opens all four paths above.",
         Inches(0.4), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=ORANGE, size=14, bold=True)

    notes(slide, """\
SLIDE 14 — Career Paths (3 min)

"Four distinct directions you can take after this course."

"Consultant: most immediate path. You have 4 portfolio projects.
Start reaching out to businesses. First client within 30 days of finishing — realistic."

"AI Solutions Engineer: companies are hiring internally now.
Job titles include 'Automation Engineer', 'AI Ops', 'RevOps Engineer' — all pay well."

"No-code product builder: the most exciting path long-term.
Build a niche SaaS automation tool — say, specifically for real estate agents.
Charge ₹5,000/month per client. 100 clients = ₹5 lakh/month."

"Marketing AI lead: combine digital marketing knowledge with automation skills.
Arguably the most in-demand hybrid role right now."

"And importantly: no prior coding required for any of these.
Python adds more options — but it's not required to start."
""")
    return slide


def s15_why_wwweb(prs):
    """Why WorldWithWeb"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "Why WorldWithWeb for This Course?")

    differentiators = [
        ("✅  We use these systems ourselves",
         "The POCs you saw are live for WorldWithWeb. You learn from practitioners, not theorists."),
        ("✅  Limited seats — personal attention",
         "Small batch size intentional. Your questions answered. Nobody falls behind."),
        ("✅  4 live builds + capstone you ship",
         "Not theory. Not slides. Week 1 you're building a live webhook. Week 4 you're shipping."),
        ("✅  Every sector, every goal",
         "Business · Marketing · Data Science/ML · Career — a track for every student's field."),
        ("✅  Freelancing module + client pitching",
         "Pricing, proposals, packaging automations as a service. The business side, not just tech."),
        ("✅  Reusable workflow library + case study",
         "You walk away with n8n templates, a documented case study, and a completion certificate."),
    ]

    for i, (heading, body) in enumerate(differentiators):
        y = Inches(1.75) + i * Inches(0.88)
        card(slide, heading, Inches(0.4), y, Inches(5.5), Inches(0.5),
             bg_c=SURFACE, fg=GREEN, size=14, bold=True)
        card(slide, body, Inches(6.1), y, Inches(6.8), Inches(0.5),
             bg_c=CARD, fg=WHITE, size=14)

    card(slide, "26+ years training IT professionals. Now building the next generation of AI automation engineers.",
         Inches(0.4), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=ORANGE, size=14)

    notes(slide, """\
SLIDE 15 — Why WorldWithWeb (2 min)

"Let me tell you why this course with us specifically."

"The most important point: we're not teaching theory.
The three POCs you just saw? Those are our real systems. We built them.
We run them. We know the edge cases and the fixes. That's what you'll learn."

"Small batches are deliberate. I'd rather run 3 batches of 10 than one of 50.
Your questions matter. Your projects get reviewed."

"The freelancing module alone is worth the fee.
Most technical courses teach you skills but leave you stranded when it comes to
finding clients, pricing, and delivering work professionally. We cover all of it."
""")
    return slide


def s16_batch_details(prs):
    """Batch details"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "Batch Details — AI Automation & Agentic AI")

    details = [
        ("Duration",       "1 Month · 16 Hands-On Sessions",                          CYAN),
        ("Fee",            "Rs. 10,000",                                              GREEN),
        ("Level",          "Beginner to Job-Ready",                                   WHITE),
        ("Mode",           "Practical + Project-Based · Recorded sessions included",  WHITE),
        ("Batch Size",     "Limited seats per batch for personal attention",           ORANGE),
        ("Prerequisites",  "Basic computer skills · No coding required",              WHITE),
        ("Tools",          "n8n · Google Sheets · Gmail · WhatsApp · Telegram · OpenAI / Gemini / Claude · AI Agent node", WHITE),
        ("Outcome",        "4 sector builds + your own capstone",                     GREEN),
        ("Certification",  "WorldWithWeb Course Certificate",                         WHITE),
        ("Batch Start",    "[DATE — confirm before session]",                         YELLOW),
    ]

    dw1, dw2 = Inches(2.8), Inches(9.0)
    dh = Inches(0.52)
    sx = Inches(0.4)
    sy = Inches(1.75)
    gap = Inches(0.14)

    for i, (label, value, color) in enumerate(details):
        y = sy + i * (dh + gap)
        card(slide, label, sx, y, dw1, dh, bg_c=SURFACE, fg=GRAY, size=13)
        card(slide, value, sx + dw1 + Inches(0.1), y, dw2, dh,
             bg_c=CARD, fg=color, size=14, bold=(color != WHITE))

    notes(slide, """\
SLIDE 16 — Batch Details (2 min)

Fill in DATE and FEE before the session.

"8 weeks. 3 sessions per week. All recorded so you never miss."

"Maximum 10 students — this is a hard cap. It's not a marketing line.
We've run larger batches and the quality drops. 10 is the number."

"No coding required to start. You'll write maybe 10 lines of JavaScript
total — and we'll walk through every line together."

"4 projects. Real systems. A portfolio you can show clients immediately."

Present fee directly — don't be vague. Professionals respect transparency.
Have an EMI option ready to mention if needed.

"We're taking registrations now. We'll close once 10 seats are filled —
usually within a week of opening."
""")
    return slide


def s17_investment_roi(prs):
    """Investment vs ROI"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "Investment & ROI", "The honest numbers")

    card(slide, "Course Investment",
         Inches(0.5), Inches(1.75), Inches(5.6), Inches(0.6),
         bg_c=SURFACE, fg=GRAY, size=16, bold=True)
    card(slide, "Rs. 10,000",
         Inches(0.5), Inches(2.4), Inches(5.6), Inches(1.2),
         bg_c=CARD, fg=CYAN, size=36, bold=True)
    card(slide, "1 month · 16 sessions · 4 builds + capstone",
         Inches(0.5), Inches(3.65), Inches(5.6), Inches(0.5),
         bg_c=SURFACE, fg=GRAY, size=13)

    card(slide, "Your First Freelance Project",
         Inches(7.2), Inches(1.75), Inches(5.6), Inches(0.6),
         bg_c=SURFACE, fg=GRAY, size=16, bold=True)
    card(slide, "₹30,000–₹60,000",
         Inches(7.2), Inches(2.4), Inches(5.6), Inches(1.2),
         bg_c=CARD, fg=GREEN, size=32, bold=True)
    card(slide, "Recovers full course investment in 1–2 projects",
         Inches(7.2), Inches(3.65), Inches(5.6), Inches(0.5),
         bg_c=SURFACE, fg=GRAY, size=13)

    txt(slide, "vs",
        Inches(5.9), Inches(2.8), Inches(1.5), Inches(0.8),
        color=ORANGE, size=28, bold=True, align=PP_ALIGN.CENTER)

    comparisons = [
        ("MBA programme", "₹5L–₹25L", RED),
        ("6-month coding bootcamp", "₹80K–₹2L", RED),
        ("Online certification (Coursera/Udemy)", "₹5K–₹30K (no mentorship)", GRAY),
        ("This course — WorldWithWeb", "Rs. 10,000 · 1 month · 4 builds + capstone · Freelancing module included", GREEN),
    ]
    for i, (item, cost, color) in enumerate(comparisons):
        y = Inches(4.4) + i * Inches(0.52)
        card(slide, item, Inches(0.5), y, Inches(6.0), Inches(0.45),
             bg_c=CARD, fg=WHITE, size=13)
        card(slide, cost, Inches(6.65), y, Inches(6.15), Inches(0.45),
             bg_c=CARD, fg=color, size=13, bold=(color == GREEN))

    notes(slide, """\
SLIDE 17 — Investment & ROI (3 min)

Fill in course fee before the session.

"Let me show you the ROI math — because professionals make data-driven decisions."

"Your first project after this course: conservatively ₹30,000–₹60,000.
That's a 1–2 project payback on the course. After that, it's pure profit."

"Compare this to other options:
MBA — ₹5 to ₹25 lakhs. Takes 2 years. No hands-on automation skills.
Bootcamp — ₹80K to ₹2L. Mostly coding, not automation business skills.
Udemy — cheap, but no mentorship, no portfolio, no structure, no accountability."

"This course is the middle path: affordable, structured, mentored, with real output."

If someone hesitates on price:
'Let me ask you — what's one hour of your time worth as a professional?
One client project pays back this entire course. The question is how quickly you want to get there.'
""")
    return slide


def s18_testimonials_social(prs):
    """Social proof"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "What Students Say", "Real outcomes from WorldWithWeb alumni")

    testimonials = [
        ("Built my first automation client project 3 weeks after the course. Charged ₹45,000. The client asked for a second one immediately.",
         "— Freelance Automation Consultant, Ludhiana", CYAN),
        ("I was running our school's admissions manually. Now an n8n agent handles every enquiry, follow-up, and fee reminder. Saved 20 hours/week.",
         "— Director, Private School, Punjab", GREEN),
        ("Added AI automation to my digital marketing services. My retainer clients now pay 40% more because I deliver automated reporting and social posting.",
         "— Digital Marketing Professional, Chandigarh", PURPLE),
    ]

    th = Inches(1.8)
    tw = Inches(12.33)
    sy = Inches(1.65)
    gap = Inches(0.3)

    for i, (quote, attr, color) in enumerate(testimonials):
        y = sy + i * (th + gap)
        card(slide, f'"{quote}"', Inches(0.4), y, tw, Inches(1.3),
             bg_c=CARD, fg=WHITE, size=14)
        card(slide, attr, Inches(0.4), y + Inches(1.32), Inches(8.0), Inches(0.4),
             bg_c=SURFACE, fg=color, size=13)

    card(slide, "1000+ students trained · 26+ years · Ludhiana, Punjab",
         Inches(0.4), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=ORANGE, size=14)

    notes(slide, """\
SLIDE 18 — Social Proof (2 min)

"Let me show you what's happened for people who've taken this path."

Read each testimonial — or better, if you have a real alum in the room, let them speak.

"These are real outcomes. Not outliers — these are the kind of results we target for every student."

If you have more testimonials / screenshots / case studies, add them here.
WhatsApp screenshots from students work very well for live sessions.
""")
    return slide


def s19_who_this_is_for(prs):
    """Who this is for / who it's not for"""
    slide = _blank(prs)
    bg(slide)
    title_block(slide, "Is This Course Right for You?")

    yes_items = [
        "Students & freshers — any field: marketing, DS/ML, web, business",
        "Working professional wanting to automate and add an income stream",
        "Business owner — stop doing repetitive work by hand",
        "Digital marketer wanting to offer AI automation services",
        "Anyone curious about AI — no coding background needed",
        "Freelancer wanting to upgrade from basic services to automation builds",
    ]
    no_items = [
        "Looking for a theory-only certificate with no practical work",
        "Not willing to practice and build between sessions",
        "Senior ML/AI researcher (this is applied, not academic)",
        "Looking for a get-rich-quick shortcut",
    ]

    txt(slide, "This IS for you if:", Inches(0.5), Inches(1.7),
        Inches(6.3), Inches(0.45), color=GREEN, size=15, bold=True)
    for i, item in enumerate(yes_items):
        card(slide, f"✅  {item}", Inches(0.5), Inches(2.2) + i * Inches(0.7),
             Inches(6.3), Inches(0.58), bg_c=CARD, fg=WHITE, size=14)

    txt(slide, "This is NOT for you if:", Inches(7.0), Inches(1.7),
        Inches(5.9), Inches(0.45), color=RED, size=15, bold=True)
    for i, item in enumerate(no_items):
        card(slide, f"✗  {item}", Inches(7.0), Inches(2.2) + i * Inches(0.7),
             Inches(5.9), Inches(0.58), bg_c=CARD, fg=GRAY, size=14)

    card(slide, "Honest qualification saves everyone's time. We want students who are ready to commit and build.",
         Inches(0.4), Inches(6.65), Inches(12.33), Inches(0.55),
         bg_c=SURFACE, fg=ORANGE, size=14)

    notes(slide, """\
SLIDE 19 — Qualification (2 min)

"I want to be upfront about who this course is designed for — and who it isn't."

"If you're a working professional, a business owner, a marketer, or someone who wants
a practical AI skill they can monetise — this is built for you."

"If you're looking for a certificate to hang on your wall without putting in the work —
this isn't the right fit. We run small batches because we want everyone to actually build."

"No coding required. We've had students from commerce, arts, business backgrounds
complete this and land clients. It's about logical thinking, not syntax."

Transition: "Any questions before I share the enrollment details?"
[Take 2–3 questions, then move to final slide]
""")
    return slide


def s20_cta(prs):
    """Final CTA"""
    slide = _blank(prs)
    bg(slide)

    txt(slide,
        "The Businesses That Automate First\nWill Dominate Their Markets.",
        Inches(0.5), Inches(0.6), Inches(12.33), Inches(1.6),
        color=CYAN, size=38, bold=True, align=PP_ALIGN.CENTER)

    sh = slide.shapes.add_shape(1, Inches(2.5), Inches(2.35), Inches(8.33), Inches(0.04))
    sh.fill.solid(); sh.fill.fore_color.rgb = CYAN; sh.line.fill.background()

    txt(slide,
        "You can build those automations.  We'll teach you exactly how.",
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(0.65),
        color=WHITE, size=22, align=PP_ALIGN.CENTER)

    card(slide,
         "Batch Start: [DATE]     ·     Fee: Rs. 10,000     ·     Limited seats",
         Inches(1.5), Inches(3.35), Inches(10.33), Inches(0.75),
         bg_c=SURFACE, fg=WHITE, size=18)

    card(slide, "Reserve Your Seat  →  Talk to us after this session",
         Inches(2.0), Inches(4.3), Inches(9.33), Inches(0.65),
         bg_c=ORANGE, fg=BG, size=18, bold=True)

    for i, (icon, contact) in enumerate([
        ("📱 WhatsApp", "[NUMBER]"),
        ("🌐 Website",  "worldwithweb.in"),
        ("📧 Email",    "[EMAIL]"),
    ]):
        x = Inches(1.5) + i * Inches(3.6)
        card(slide, icon, x, Inches(5.2), Inches(1.4), Inches(0.5),
             bg_c=CARD, fg=GRAY, size=13)
        card(slide, contact, x + Inches(1.45), Inches(5.2), Inches(2.1), Inches(0.5),
             bg_c=CARD, fg=CYAN, size=13, bold=True)

    txt(slide, "WorldWithWeb  ·  Ludhiana's First AI Automation & Agentic AI Course  ·  Build, don't just watch.",
        Inches(0.5), Inches(6.75), Inches(12.33), Inches(0.45),
        color=ORANGE, size=13, bold=True, align=PP_ALIGN.CENTER)

    notes(slide, """\
SLIDE 20 — Final CTA (5 min)

[Read headline slowly]
"The businesses that automate first will dominate their markets."

[Pause]

"You've just seen three live automations that are running for WorldWithWeb right now.
You've seen what the market pays. You've seen the career paths."

"The question isn't whether AI automation is the future — it's whether YOU want
to be the one building it or the one hiring someone who does."

[Point to enrollment box]
"Batch starts [DATE]. Maximum 10 students. Fee is [AMOUNT]."
"I'm right here after this. Come talk to me. No script, no pressure."
"If this course is right for you, we'll know in 5 minutes."

POST-SESSION:
- Collect name + WhatsApp number + course interest
- Send a WhatsApp confirmation to everyone who shows interest within 2 hours
- Follow up within 24 hours with full course brochure + payment link
- Log all leads in the enquiry sheet
""")
    return slide


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    builders = [
        s01_hook, s02_chatbot_vs_agent, s03_what_agents_do, s04_market_opportunity,
        s05_poc1, s06_poc2, s07_poc3, s08_usecase_business,
        s09_usecase_marketer, s10_what_youll_build, s11_tools, s12_course_structure,
        s13_freelance_income, s14_career_paths, s15_why_wwweb, s16_batch_details,
        s17_investment_roi, s18_testimonials_social, s19_who_this_is_for, s20_cta,
    ]

    for i, fn in enumerate(builders, 1):
        print(f"  Building slide {i:02d}/20 — {fn.__doc__}")
        fn(prs)

    out_path = os.path.join(
        "/Users/macsolutions/Projects/worldwithweb_assets/decks/session_decks",
        "WorldWithWeb_AI_Agents_Course_Deck.pptx"
    )
    prs.save(out_path)
    print(f"\n✅  Saved → {out_path}")


if __name__ == "__main__":
    main()
