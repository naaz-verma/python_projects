"""
Generator for WorldWithWeb_AI_Motivation_Session.pptx
Run: python sessions/create_ppt_ai_motivation_session.py
Output: presentations/WorldWithWeb_AI_Motivation_Session.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
SURFACE  = RGBColor(0x25, 0x25, 0x40)   # card bg
CYAN     = RGBColor(0x00, 0xB4, 0xD8)   # primary accent / titles
GREEN    = RGBColor(0x00, 0xE6, 0x76)   # highlights / yes
PURPLE   = RGBColor(0xBB, 0x86, 0xFC)   # secondary accent
RED      = RGBColor(0xFF, 0x45, 0x6E)   # no / warning
ORANGE   = RGBColor(0xFF, 0x9F, 0x1C)   # CTAs / progression
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x99, 0x99, 0x99)
DARK_BG2 = RGBColor(0x0D, 0x11, 0x17)   # deeper dark for flow boxes

# ── Slide dimensions (widescreen 13.33" × 7.5") ─────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ── Low-level helpers ────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text: str, left, top, width, height,
                color: RGBColor = WHITE,
                size: int = 24,
                bold: bool = False,
                align=PP_ALIGN.LEFT,
                italic: bool = False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_card(slide, text: str, left, top, width, height,
             bg_color: RGBColor = SURFACE,
             text_color: RGBColor = WHITE,
             text_size: int = 18,
             bold: bool = False,
             radius: bool = True):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()  # no border

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(text_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    return shape


def add_title(slide, text: str, subtitle: str = ""):
    add_textbox(slide, text,
                left=Inches(0.5), top=Inches(0.25),
                width=Inches(12.33), height=Inches(0.85),
                color=CYAN, size=36, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    left=Inches(0.5), top=Inches(1.05),
                    width=Inches(12.33), height=Inches(0.5),
                    color=GRAY, size=16, align=PP_ALIGN.LEFT)


def add_divider(slide, top=Inches(1.1)):
    shape = slide.shapes.add_shape(1,
        Inches(0.5), top, Inches(12.33), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()


def set_speaker_notes(slide, notes: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes


# ── Slide builders ────────────────────────────────────────────────────────────

def build_slide_01(prs):
    """The Hook"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)

    # Faint watermark chips (tech logos as text)
    watermarks = ["ChatGPT", "Tesla", "Midjourney", "OpenAI", "Gemini", "Claude"]
    positions = [
        (Inches(0.3), Inches(0.4)),
        (Inches(9.5), Inches(0.3)),
        (Inches(1.5), Inches(6.5)),
        (Inches(10.2), Inches(6.2)),
        (Inches(0.2), Inches(3.5)),
        (Inches(10.8), Inches(3.2)),
    ]
    for wm, (lft, tp) in zip(watermarks, positions):
        add_textbox(slide, wm, lft, tp, Inches(2.5), Inches(0.45),
                    color=RGBColor(0x33, 0x33, 0x55), size=18, bold=True)

    # Main headline
    add_textbox(
        slide,
        "The Next Billion-Dollar Companies\nWill Be Built With AI.",
        left=Inches(1), top=Inches(1.6),
        width=Inches(11.33), height=Inches(2),
        color=CYAN, size=44, bold=True, align=PP_ALIGN.CENTER
    )

    # Sub-headline
    add_textbox(
        slide,
        "Will You Build Them  —  or Work For Them?",
        left=Inches(1), top=Inches(3.7),
        width=Inches(11.33), height=Inches(0.8),
        color=WHITE, size=28, bold=False, align=PP_ALIGN.CENTER
    )

    # Brand tag
    add_textbox(
        slide, "WorldWithWeb",
        left=Inches(5.16), top=Inches(6.7),
        width=Inches(3), height=Inches(0.5),
        color=ORANGE, size=14, bold=True, align=PP_ALIGN.CENTER
    )

    set_speaker_notes(slide, """\
OPENING — The Hook (3 min)

Walk in confidently. Let the slide sit for 5 seconds. Then:

ASK: "Quick show of hands — who here has used ChatGPT?"
[Almost every hand goes up]

"Great. Now keep your hand up if you actually know HOW ChatGPT works."
[Most hands go down]

[Pause. Let that land.]

"That gap you just saw — that's the entire career opportunity we're talking about today.
Consumers use AI. Professionals BUILD with AI."

"One year from now, every company — healthcare, banking, retail, entertainment —
will run on AI. The question is simple: are you going to be the one using it,
or the one who BUILT it?"

Transition: "Let me show you what's really happening out there first."
""")
    return slide


def build_slide_02(prs):
    """Reality Check — Jobs Are Changing"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Jobs Are Changing Faster Than Degrees")
    add_divider(slide)

    cols = [
        ("2010", ["Typist", "Data Entry Clerk", "Telephone Operator", "Fax Operator"], SURFACE, GRAY),
        ("2020", ["Social Media Manager", "App Developer", "SEO Specialist", "Cloud Engineer"], SURFACE, WHITE),
        ("2030", ["AI Engineer", "Prompt Engineer", "AI Agent Developer", "AI Automation\nConsultant", "AI Product Manager"], CYAN, BG),
    ]
    col_w = Inches(3.8)
    col_gap = Inches(0.35)
    start_x = Inches(0.75)

    for i, (year, jobs, bg_c, fg_c) in enumerate(cols):
        x = start_x + i * (col_w + col_gap)
        # Year header
        add_card(slide, year, x, Inches(1.4), col_w, Inches(0.55),
                 bg_color=ORANGE if year == "2030" else SURFACE,
                 text_color=BG if year == "2030" else CYAN,
                 text_size=22, bold=True)
        # Job chips
        for j, job in enumerate(jobs):
            add_card(slide, job, x, Inches(2.1) + j * Inches(0.95),
                     col_w, Inches(0.82),
                     bg_color=bg_c, text_color=fg_c,
                     text_size=15, bold=(year == "2030"))

    set_speaker_notes(slide, """\
SLIDE 2 — Reality Check (3 min)

"Look at this. In 2010, companies were hiring for typing and data entry.
By 2020, you needed to know apps and social media.
In 2030 — which is only 4 years away — the most in-demand roles are ALL AI roles."

Point to the 2030 column:
"AI Engineer. Prompt Engineer. AI Automation Consultant. AI Agent Developer.
These jobs DIDN'T EXIST 5 years ago. And they're already paying more than traditional IT."

Key message:
"Here's the truth: AI won't replace everyone. But people using AI WILL replace
people who don't. The choice isn't AI vs humans — it's AI-skilled humans vs unskilled ones."
""")
    return slide


def build_slide_03(prs):
    """AI Is Already Everywhere"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "AI Is Already in Your Pocket")
    add_divider(slide)

    brands = [
        ("Netflix", CYAN),
        ("Instagram", PURPLE),
        ("YouTube", RED),
        ("Amazon", ORANGE),
        ("Uber", GREEN),
        ("Google Maps", CYAN),
        ("Spotify", GREEN),
    ]

    positions = [
        (Inches(0.6),  Inches(1.9)),
        (Inches(3.7),  Inches(1.9)),
        (Inches(6.8),  Inches(1.9)),
        (Inches(9.9),  Inches(1.9)),
        (Inches(1.4),  Inches(4.1)),
        (Inches(4.7),  Inches(4.1)),
        (Inches(8.0),  Inches(4.1)),
    ]
    chip_w, chip_h = Inches(2.6), Inches(1.0)

    for (brand, color), (lft, tp) in zip(brands, positions):
        add_card(slide, brand, lft, tp, chip_w, chip_h,
                 bg_color=SURFACE, text_color=color,
                 text_size=20, bold=True)

    add_textbox(
        slide,
        "How does Netflix know what you want to watch?   How does Amazon know what to recommend?",
        left=Inches(0.5), top=Inches(5.7),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=15, align=PP_ALIGN.CENTER, italic=True
    )
    add_textbox(
        slide, "AI.",
        left=Inches(5.9), top=Inches(6.3),
        width=Inches(1.5), height=Inches(0.6),
        color=CYAN, size=32, bold=True, align=PP_ALIGN.CENTER
    )

    set_speaker_notes(slide, """\
SLIDE 3 — AI Is Already Everywhere (3 min)

"Before we talk about learning AI, let me show you something."

Point to Netflix: "How does Netflix know exactly what to recommend after you finish a show?"
[Let someone answer]

"How does Instagram decide which posts show up first? How does Spotify build your
Discover Weekly playlist? How does Google Maps know there's traffic ahead?"

"Every single one of these is AI. You use it 50 times a day without thinking about it."

"These companies — Netflix, Amazon, Google — employ thousands of AI engineers.
Each one started exactly where you are."

Transition: "And this is just the beginning. The gold rush is happening RIGHT NOW."
""")
    return slide


def build_slide_04(prs):
    """The AI Gold Rush"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Every Company Is Now an AI Company")
    add_divider(slide)

    companies = [
        ("OpenAI", "$157B+", CYAN),
        ("Anthropic", "$61B+", PURPLE),
        ("Google\nGemini", "$2T+", GREEN),
        ("Nvidia", "$3T+", ORANGE),
    ]
    cw, ch = Inches(2.8), Inches(1.6)
    cx = Inches(0.55)
    for i, (name, val, color) in enumerate(companies):
        x = cx + i * (cw + Inches(0.3))
        add_card(slide, name, x, Inches(1.5), cw, ch,
                 bg_color=SURFACE, text_color=color, text_size=18, bold=True)
        add_textbox(slide, val, x, Inches(3.15), cw, Inches(0.45),
                    color=WHITE, size=14, align=PP_ALIGN.CENTER)

    sectors = ["Healthcare", "Banking", "Retail", "Marketing", "Manufacturing", "Education", "Legal"]
    add_textbox(slide, "Every sector is hiring AI talent right now:",
                left=Inches(0.5), top=Inches(4.0),
                width=Inches(12.33), height=Inches(0.4),
                color=GRAY, size=15)

    sw, sh = Inches(1.65), Inches(0.65)
    for i, sec in enumerate(sectors):
        row = i // 4
        col = i % 4
        add_card(slide, sec,
                 Inches(0.5) + col * (sw + Inches(0.2)),
                 Inches(4.55) + row * (sh + Inches(0.15)),
                 sw, sh,
                 bg_color=SURFACE, text_color=WHITE, text_size=14)

    set_speaker_notes(slide, """\
SLIDE 4 — The AI Gold Rush (3 min)

"Look at these numbers. OpenAI — over $150 billion. Nvidia — over $3 TRILLION.
These aren't just tech companies any more. They are the new oil."

"And the demand isn't just in Silicon Valley. Healthcare companies need AI to read scans.
Banks need AI for fraud detection. Retailers need AI for recommendations.
Every sector you can think of — they all need people who understand AI."

"India alone is projected to have 1 million AI job openings by 2026.
The talent supply isn't keeping up. That is the gap you can fill."

Transition: "Let me show you what AI can actually DO — live, right now."
""")
    return slide


def build_slide_05(prs):
    """Demo 1 — AI Image Generation"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "DEMO #1 — AI Image Generation", subtitle="Live demo")
    add_divider(slide)

    # Prompt box
    add_card(slide, "PROMPT",
             Inches(0.5), Inches(1.5), Inches(2.0), Inches(0.5),
             bg_color=ORANGE, text_color=BG, text_size=13, bold=True)
    add_card(slide,
             '"Create a Nike advertisement with a futuristic athlete running on neon streets"',
             Inches(0.5), Inches(2.05), Inches(5.8), Inches(1.2),
             bg_color=SURFACE, text_color=WHITE, text_size=16)

    # Arrow
    add_textbox(slide, "→",
                Inches(6.5), Inches(2.35), Inches(0.8), Inches(0.6),
                color=CYAN, size=36, bold=True, align=PP_ALIGN.CENTER)

    # Result placeholder
    add_card(slide, "🎨  AI-Generated\nNike Ad\n[Show live result here]",
             Inches(7.4), Inches(1.5), Inches(5.4), Inches(3.5),
             bg_color=SURFACE, text_color=CYAN, text_size=18, bold=True)

    # Time comparison
    add_card(slide, "Designer: 2-3 hours",
             Inches(0.5), Inches(4.8), Inches(4.0), Inches(0.7),
             bg_color=SURFACE, text_color=RED, text_size=16)
    add_card(slide, "AI: 8 seconds",
             Inches(4.8), Inches(4.8), Inches(4.0), Inches(0.7),
             bg_color=SURFACE, text_color=GREEN, text_size=16, bold=True)

    add_textbox(
        slide,
        "AI didn't replace the designer. It gave superpowers to the one who knows how to use it.",
        left=Inches(0.5), top=Inches(5.7),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=14, align=PP_ALIGN.CENTER, italic=True
    )

    set_speaker_notes(slide, """\
SLIDE 5 — Demo #1: AI Image Generation (3-4 min)

[Open Midjourney / DALL-E / Adobe Firefly in browser]

"I'm going to type one sentence. Watch what happens."

Type the prompt live: "Create a Nike advertisement with a futuristic athlete running on neon streets"

[Show result — pause for reactions]

"How long would a professional designer take for this? At least 2-3 hours.
AI did it in 8 seconds."

"Now here's the key insight: the AI didn't think of that idea. I did.
The creativity, the brief, the direction — that still comes from a human.
AI is the execution engine. The human is still the strategist."

"This is the new creative workflow. And companies are hiring people who know it."
""")
    return slide


def build_slide_06(prs):
    """Demo 2 — AI Website Builder"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "DEMO #2 — AI Builds a Website", subtitle="Live demo")
    add_divider(slide)

    add_card(slide, "PROMPT",
             Inches(0.5), Inches(1.5), Inches(2.0), Inches(0.5),
             bg_color=ORANGE, text_color=BG, text_size=13, bold=True)
    add_card(slide,
             '"Create a modern website for a premium coffee shop called BrewLab. Include a menu, about section, and contact form."',
             Inches(0.5), Inches(2.05), Inches(5.8), Inches(1.4),
             bg_color=SURFACE, text_color=WHITE, text_size=16)

    add_textbox(slide, "→",
                Inches(6.5), Inches(2.5), Inches(0.8), Inches(0.6),
                color=CYAN, size=36, bold=True, align=PP_ALIGN.CENTER)

    add_card(slide, "🌐  Full Website\nHTML + CSS + JS\n[Show live result here]",
             Inches(7.4), Inches(1.5), Inches(5.4), Inches(3.5),
             bg_color=SURFACE, text_color=CYAN, text_size=18, bold=True)

    add_textbox(
        slide,
        "Web development is not dying — it's evolving. The builders who direct AI will be paid more than those who just write code.",
        left=Inches(0.5), top=Inches(5.7),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=14, align=PP_ALIGN.CENTER, italic=True
    )

    set_speaker_notes(slide, """\
SLIDE 6 — Demo #2: AI Website Builder (3 min)

[Open v0.dev / Cursor / Bolt.new or similar AI web builder]

"Now let's build a website. Whole website. In 60 seconds."

[Type prompt and show result live]

"Full HTML, CSS, responsive design — done.

A freelance developer would charge ₹15,000–₹50,000 for this.
A student who knows how to direct AI tools can deliver this in a client meeting.

This is NOT the end of web development. It's the beginning of AI-assisted development.
The programmers who survive — and thrive — are the ones who understand what the AI is doing
well enough to direct it, fix it, and improve it."

Transition: "Let me show you something even wilder."
""")
    return slide


def build_slide_07(prs):
    """Demo 3 — AI Voice Clone"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "DEMO #3 — AI Voice Clone", subtitle="Live demo")
    add_divider(slide)

    add_card(slide, "YOUR VOICE",
             Inches(1.0), Inches(2.2), Inches(4.5), Inches(0.7),
             bg_color=SURFACE, text_color=CYAN, text_size=18, bold=True)

    # Waveform simulation (bars of varying height)
    wave_heights = [0.3, 0.6, 0.9, 0.5, 1.1, 0.4, 0.8, 0.35, 0.7, 0.55, 1.0, 0.45]
    bar_w = Inches(0.28)
    for i, h in enumerate(wave_heights):
        bar_h = Inches(h)
        bar_top = Inches(3.35) - bar_h / 2
        shape = slide.shapes.add_shape(1,
            Inches(1.1) + i * Inches(0.3), bar_top, bar_w, bar_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CYAN
        shape.line.fill.background()

    add_card(slide, "AI CLONED VOICE",
             Inches(7.5), Inches(2.2), Inches(4.5), Inches(0.7),
             bg_color=SURFACE, text_color=PURPLE, text_size=18, bold=True)

    for i, h in enumerate(wave_heights):
        bar_h = Inches(h)
        bar_top = Inches(3.35) - bar_h / 2
        shape = slide.shapes.add_shape(1,
            Inches(7.6) + i * Inches(0.3), bar_top, bar_w, bar_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PURPLE
        shape.line.fill.background()

    add_textbox(slide, "vs",
                Inches(6.1), Inches(3.0), Inches(1.1), Inches(0.7),
                color=ORANGE, size=28, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(
        slide,
        "AI can replicate any voice with just 10 seconds of audio.\nUnderstanding this is as important as being excited by it.",
        left=Inches(0.5), top=Inches(5.5),
        width=Inches(12.33), height=Inches(0.8),
        color=GRAY, size=15, align=PP_ALIGN.CENTER, italic=True
    )

    set_speaker_notes(slide, """\
SLIDE 7 — Demo #3: AI Voice Clone (3 min)

[Open ElevenLabs or similar voice cloning tool]

"This one always gets a reaction."

[Record or play your own voice — then play the AI clone]

"Same voice. Different source. The AI learned from 10 seconds of audio."

[Pause for audience reaction — let them process it]

"Now — this is powerful. And like all powerful tools, it has a good side and a concerning side.
The good side: AI voices are used in audiobooks, accessibility tools, language learning,
customer service. The concerning side is why we need people who UNDERSTAND how this works —
not just people who use it without thinking."

"Every powerful technology requires educated, ethical builders. That's what we're training."

Transition: "Now let me show you the demo that companies are actually paying for today."
""")
    return slide


def build_slide_08(prs):
    """Demo 4 — AI Agent"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "DEMO #4 — AI Agent in Action", subtitle="No human intervention required")
    add_divider(slide)

    steps = [
        ("Student enquiry arrives on website", CYAN),
        ("AI replies instantly — answers questions", PURPLE),
        ("AI qualifies the lead — course & budget fit", ORANGE),
        ("AI sends personalised WhatsApp message", GREEN),
        ("AI updates CRM / Excel with student data", CYAN),
        ("AI generates daily summary report", ORANGE),
    ]

    box_w = Inches(8.0)
    box_h = Inches(0.65)
    start_x = Inches(2.6)
    start_y = Inches(1.45)
    gap = Inches(0.82)

    for i, (text, color) in enumerate(steps):
        y = start_y + i * gap
        add_card(slide, text, start_x, y, box_w, box_h,
                 bg_color=SURFACE, text_color=color, text_size=16, bold=(i == 0))
        if i < len(steps) - 1:
            add_textbox(slide, "↓",
                        start_x + Inches(3.5), y + box_h,
                        Inches(1.0), Inches(0.22),
                        color=GRAY, size=14, align=PP_ALIGN.CENTER)

    add_textbox(
        slide, "Companies are paying ₹30,000–₹2,00,000/month for this skill.",
        left=Inches(0.5), top=Inches(6.8),
        width=Inches(12.33), height=Inches(0.45),
        color=ORANGE, size=15, bold=True, align=PP_ALIGN.CENTER
    )

    set_speaker_notes(slide, """\
SLIDE 8 — Demo #4: AI Agent (5 min — STRONGEST DEMO)

This is your most important demo. Take your time.

[If possible, show a live n8n / Make / custom Python agent workflow]

"This is an AI agent. It's not a chatbot. It doesn't just answer questions —
it ACTS. It connects to your WhatsApp, your CRM, your spreadsheet, your calendar.
All automatically."

Walk through each step:
"A student fills in a form on the website at 2am.
The AI replies within 3 seconds. Answers questions about the course, fee, schedule.
Qualifies whether this student is a good fit.
If yes — sends a WhatsApp message automatically.
Updates the admissions spreadsheet.
Generates a morning report for the counsellor."

"Zero human involvement. The counsellor comes in the morning and sees qualified, warm leads.
Not cold calls."

"Businesses are actively paying for this skill RIGHT NOW. Not in 5 years. Today.
This is Level 6 of what you'll learn at WorldWithWeb."

[Pause for effect]

"Now — can anyone learn this? Let me show you."
""")
    return slide


def build_slide_09(prs):
    """Can Anyone Learn AI?"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Can Anyone Learn AI?")
    add_divider(slide)

    no_items = [
        "❌  Only engineers",
        "❌  Only coders",
        "❌  Only math toppers",
        "❌  Only people with laptops",
    ]
    yes_items = [
        "✅  Students (any stream)",
        "✅  Graduates & freshers",
        "✅  Working professionals",
        "✅  Business owners",
    ]

    # NO column
    add_textbox(slide, "MYTH", left=Inches(1.0), top=Inches(1.5),
                width=Inches(5.0), height=Inches(0.5),
                color=RED, size=16, bold=True, align=PP_ALIGN.CENTER)
    for i, item in enumerate(no_items):
        add_card(slide, item,
                 Inches(1.0), Inches(2.1) + i * Inches(1.0),
                 Inches(5.0), Inches(0.82),
                 bg_color=SURFACE, text_color=RED, text_size=17)

    # YES column
    add_textbox(slide, "REALITY", left=Inches(7.3), top=Inches(1.5),
                width=Inches(5.0), height=Inches(0.5),
                color=GREEN, size=16, bold=True, align=PP_ALIGN.CENTER)
    for i, item in enumerate(yes_items):
        add_card(slide, item,
                 Inches(7.3), Inches(2.1) + i * Inches(1.0),
                 Inches(5.0), Inches(0.82),
                 bg_color=SURFACE, text_color=GREEN, text_size=17)

    add_textbox(
        slide,
        "You need curiosity, consistency, and the right guidance. Not a CS degree.",
        left=Inches(0.5), top=Inches(6.6),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=15, align=PP_ALIGN.CENTER, italic=True
    )

    set_speaker_notes(slide, """\
SLIDE 9 — Can Anyone Learn AI? (2 min)

"I get this question every session. 'Is AI only for engineers? Do I need to be a programmer?'

Here's the truth:"

[Point to the MYTH column]
"AI is NOT only for engineers. NOT only for people who already code.
NOT only for toppers."

[Point to the REALITY column]
"Students from any stream — arts, commerce, science — can learn this.
Freshers can learn this. Working professionals can upskill.
Business owners can automate their entire operations."

"The only requirement is curiosity and consistency.
We take care of the structure, the projects, and the mentorship."

Transition: "So what does the actual learning path look like? Let me show you."
""")
    return slide


def build_slide_10(prs):
    """Your Journey — 8-Level Roadmap"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Your Journey")
    add_divider(slide)

    levels = [
        ("L1", "Python",              ORANGE),
        ("L2", "Data Analysis",       ORANGE),
        ("L3", "Machine Learning",    CYAN),
        ("L4", "Artificial Intelligence", CYAN),
        ("L5", "Generative AI",       PURPLE),
        ("L6", "AI Agents",           PURPLE),
        ("L7", "Industry Projects",   GREEN),
        ("L8", "Job Ready  🎯",        GREEN),
    ]

    box_w = Inches(2.7)
    box_h = Inches(0.58)
    cols = 4
    col_gap = Inches(0.25)
    row_gap = Inches(0.3)
    start_x = Inches(0.55)
    start_y = Inches(1.55)

    for i, (lbl, name, color) in enumerate(levels):
        row = i // cols
        col = i % cols
        x = start_x + col * (box_w + col_gap)
        y = start_y + row * (box_h + row_gap)
        # level badge
        add_card(slide, lbl, x, y, Inches(0.55), box_h,
                 bg_color=color, text_color=BG, text_size=13, bold=True)
        # name
        add_card(slide, name, x + Inches(0.6), y, Inches(2.1), box_h,
                 bg_color=SURFACE, text_color=WHITE, text_size=14)
        # arrow (except last in row and last overall)
        if col < cols - 1 and i < len(levels) - 1:
            add_textbox(slide, "→",
                        x + box_w + Inches(0.02), y,
                        Inches(0.2), box_h,
                        color=GRAY, size=12, align=PP_ALIGN.CENTER)

    add_textbox(
        slide,
        "Think of it like a game. Every level unlocks the next. You don't skip levels — you master them.",
        left=Inches(0.5), top=Inches(5.3),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=15, align=PP_ALIGN.CENTER, italic=True
    )

    add_card(slide, "Complete roadmap: 6–9 months  |  Part-time friendly  |  Project at every level",
             Inches(0.5), Inches(5.95), Inches(12.33), Inches(0.7),
             bg_color=SURFACE, text_color=ORANGE, text_size=15)

    set_speaker_notes(slide, """\
SLIDE 10 — Your Journey (4 min)

"Let me show you the full roadmap. This is what your learning journey looks like."

Walk through each level:
"Level 1: Python. The language of AI. Everything starts here.
Level 2: Data Analysis. Because AI runs on data — you need to understand it.
Level 3: Machine Learning. Teaching computers to learn patterns.
Level 4: AI. Computer vision, natural language processing.
Level 5: Generative AI. Building with ChatGPT, Claude, Gemini.
Level 6: AI Agents. Autonomous systems that act without humans.
Level 7: Industry Projects. Real client-grade work in your portfolio.
Level 8: Job Ready. Interview prep, portfolio, LinkedIn, referrals."

"Think of this like a video game. You can't skip to Level 6 without the foundation.
But every single level builds directly on the last. There's no wasted time."

"This entire journey can be done part-time — alongside college or work.
And we have a project at every single level, so you're building from Day 1."
""")
    return slide


def build_slide_11(prs):
    """Level 1: Python"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Level 1: Python — The Language of AI")
    add_divider(slide)

    add_card(slide, "L1", Inches(0.5), Inches(1.45), Inches(0.6), Inches(0.5),
             bg_color=ORANGE, text_color=BG, text_size=14, bold=True)
    add_textbox(slide, "The entry point. Every AI engineer in the world started here.",
                left=Inches(1.2), top=Inches(1.45),
                width=Inches(11.0), height=Inches(0.5),
                color=GRAY, size=15)

    projects = [
        ("Calculator & Logic", "Variables, loops, functions — core thinking tools"),
        ("Automation Scripts", "Automate repetitive tasks — email, files, data"),
        ("Data Processing", "Read, clean, and analyse real-world datasets"),
        ("Chatbot (Rule-based)", "Your first conversational AI — no ML needed yet"),
    ]
    pw = Inches(5.9)
    ph = Inches(1.2)
    gap = Inches(0.25)

    for i, (title, desc) in enumerate(projects):
        row = i // 2
        col = i % 2
        x = Inches(0.5) + col * (pw + gap)
        y = Inches(2.2) + row * (ph + gap)
        add_card(slide, title, x, y, pw, Inches(0.55),
                 bg_color=ORANGE, text_color=BG, text_size=16, bold=True)
        add_card(slide, desc, x, y + Inches(0.55), pw, Inches(0.6),
                 bg_color=SURFACE, text_color=WHITE, text_size=14)

    add_textbox(
        slide,
        "No prior experience needed. If you can use a smartphone, you can learn Python.",
        left=Inches(0.5), top=Inches(6.55),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=15, align=PP_ALIGN.CENTER, italic=True
    )

    set_speaker_notes(slide, """\
SLIDE 11 — Python (2 min)

"Every journey starts with Python. It's not just a programming language —
it's the language that ALL of AI is built in.
GPT. Stable Diffusion. TensorFlow. All Python under the hood."

"In the Python module, you'll build four projects:
A calculator — sounds simple, but it teaches you how to THINK like a programmer.
Automation scripts — imagine writing code that does your boring Excel work for you.
Data processing — you'll work with real data from companies.
And your first chatbot — no machine learning needed, pure logic."

"No experience needed at all. We start from zero. If you can send a WhatsApp, you can learn Python."
""")
    return slide


def build_slide_12(prs):
    """Level 2: Data Science"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Level 2: Data Science — The Fuel of AI")
    add_divider(slide)

    add_card(slide, "L2", Inches(0.5), Inches(1.45), Inches(0.6), Inches(0.5),
             bg_color=ORANGE, text_color=BG, text_size=14, bold=True)
    add_textbox(slide, "Data is what AI learns from. Without data skills, you can't build real AI.",
                left=Inches(1.2), top=Inches(1.45),
                width=Inches(11.0), height=Inches(0.5),
                color=GRAY, size=15)

    add_card(slide, "📊  Interactive\nBusiness Dashboard\n[Screenshot placeholder]",
             Inches(0.5), Inches(2.1), Inches(5.0), Inches(3.8),
             bg_color=SURFACE, text_color=CYAN, text_size=17, bold=True)

    projects_r = [
        ("Sales Analysis", "Find trends, peaks, drops in real sales data", CYAN),
        ("Customer Segmentation", "Group customers by behaviour using clustering", PURPLE),
        ("Business Dashboard", "Interactive charts with Plotly & Streamlit", ORANGE),
    ]
    for i, (title, desc, color) in enumerate(projects_r):
        x = Inches(5.8)
        y = Inches(2.1) + i * Inches(1.35)
        add_card(slide, title, x, y, Inches(7.0), Inches(0.55),
                 bg_color=color, text_color=BG, text_size=16, bold=True)
        add_card(slide, desc, x, y + Inches(0.55), Inches(7.0), Inches(0.7),
                 bg_color=SURFACE, text_color=WHITE, text_size=14)

    set_speaker_notes(slide, """\
SLIDE 12 — Data Science (2 min)

"Data Science is where Python meets the real world.
Every AI model learns from data. If you don't understand data, you can't build AI."

"In this module, you'll analyse real business datasets:
Sales data — finding which products are growing, which are dying.
Customer data — grouping people by behaviour, predicting churn.
And you'll build an interactive dashboard that any business owner would actually use."

"By the end, you'll be comfortable with Pandas, Matplotlib, Plotly — the standard tools
every data team uses."
""")
    return slide


def build_slide_13(prs):
    """Level 3: Machine Learning"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Level 3: Machine Learning")
    add_divider(slide)

    add_card(slide, "L3", Inches(0.5), Inches(1.45), Inches(0.6), Inches(0.5),
             bg_color=CYAN, text_color=BG, text_size=14, bold=True)
    add_textbox(slide, "Teaching computers to learn patterns — without being explicitly programmed.",
                left=Inches(1.2), top=Inches(1.45),
                width=Inches(11.0), height=Inches(0.5),
                color=GRAY, size=15)

    projects = [
        ("🏠  House Price Prediction", "Predict property values from 20+ features", CYAN),
        ("👥  Employee Attrition", "Predict who will quit — before they do", PURPLE),
        ("🎬  Movie Recommendation", "Build the Netflix engine — from scratch", ORANGE),
    ]
    pw, ph = Inches(3.8), Inches(2.8)
    gap = Inches(0.27)
    for i, (title, desc, color) in enumerate(projects):
        x = Inches(0.5) + i * (pw + gap)
        add_card(slide, title, x, Inches(2.1), pw, Inches(0.65),
                 bg_color=color, text_color=BG, text_size=15, bold=True)
        add_card(slide, desc, x, Inches(2.75), pw, Inches(0.75),
                 bg_color=SURFACE, text_color=WHITE, text_size=14)
        add_card(slide, "Scikit-learn  ·  Pandas  ·  Matplotlib",
                 x, Inches(3.5), pw, Inches(0.5),
                 bg_color=DARK_BG2, text_color=GRAY, text_size=12)

    add_textbox(
        slide,
        "Netflix, Spotify, Amazon — their recommendation engines are ML at heart. You'll build one.",
        left=Inches(0.5), top=Inches(5.7),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=15, align=PP_ALIGN.CENTER, italic=True
    )

    set_speaker_notes(slide, """\
SLIDE 13 — Machine Learning (2 min)

"Machine learning is where things get exciting. You're not writing rules any more —
you're teaching the computer to find patterns in data by itself."

"Three projects here:
House price prediction — a classic ML problem. You'll use real estate data.
Employee attrition — companies actually PAY for this. Predict who's going to leave.
Movie recommendation — this is the Netflix algorithm. You'll build it yourself."

"You mentioned Netflix on Slide 3 — well, here's where you actually BUILD that system."

"Tools: Scikit-learn, Pandas. Industry standard. Every ML job uses these."
""")
    return slide


def build_slide_14(prs):
    """Level 4: Artificial Intelligence"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Level 4: Artificial Intelligence")
    add_divider(slide)

    add_card(slide, "L4", Inches(0.5), Inches(1.45), Inches(0.6), Inches(0.5),
             bg_color=CYAN, text_color=BG, text_size=14, bold=True)
    add_textbox(slide, "Computer vision, natural language processing — AI that sees and understands.",
                left=Inches(1.2), top=Inches(1.45),
                width=Inches(11.0), height=Inches(0.5),
                color=GRAY, size=15)

    projects = [
        ("👁️  Computer Vision", "Teach AI to identify objects in images", CYAN),
        ("💬  NLP", "Sentiment analysis, text classification", PURPLE),
        ("🔍  Face Detection", "Real-time face recognition system", GREEN),
        ("📄  OCR System", "Extract text from any image or PDF", ORANGE),
    ]
    pw, ph = Inches(2.9), Inches(2.5)
    gap = Inches(0.26)
    for i, (title, desc, color) in enumerate(projects):
        x = Inches(0.5) + i * (pw + gap)
        add_card(slide, title, x, Inches(2.1), pw, Inches(0.65),
                 bg_color=color, text_color=BG, text_size=15, bold=True)
        add_card(slide, desc, x, Inches(2.75), pw, Inches(0.8),
                 bg_color=SURFACE, text_color=WHITE, text_size=14)

    add_textbox(
        slide,
        "This is where you move from predicting numbers to understanding the real world — images, language, documents.",
        left=Inches(0.5), top=Inches(5.7),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=15, align=PP_ALIGN.CENTER, italic=True
    )

    set_speaker_notes(slide, """\
SLIDE 14 — Artificial Intelligence (2 min)

"Now we're into real AI. This is where you build systems that can SEE and READ."

"Computer Vision: you'll build a model that can identify objects in photos.
Same technology self-driving cars use — at a beginner level.
NLP: teaching AI to understand human language — sentiment, meaning, intent.
Face Detection: real-time face recognition using a webcam. Students love this one.
OCR: extract text from any image or PDF automatically — massive for document processing."

"These are production-grade AI techniques. Used in healthcare, security, banking."
""")
    return slide


def build_slide_15(prs):
    """Level 5: Generative AI"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Level 5: Generative AI")
    add_divider(slide)

    add_card(slide, "L5", Inches(0.5), Inches(1.45), Inches(0.6), Inches(0.5),
             bg_color=PURPLE, text_color=WHITE, text_size=14, bold=True)
    add_textbox(slide, "Stop USING these models. Start BUILDING with them.",
                left=Inches(1.2), top=Inches(1.45),
                width=Inches(11.0), height=Inches(0.5),
                color=GRAY, size=15)

    models = [
        ("ChatGPT\nOpenAI", CYAN),
        ("Claude\nAnthropic", PURPLE),
        ("Gemini\nGoogle", GREEN),
        ("Midjourney\nImage AI", ORANGE),
    ]
    mw, mh = Inches(2.8), Inches(1.5)
    for i, (name, color) in enumerate(models):
        x = Inches(0.55) + i * (mw + Inches(0.3))
        add_card(slide, name, x, Inches(2.1), mw, mh,
                 bg_color=SURFACE, text_color=color, text_size=20, bold=True)

    skills = [
        "Prompt Engineering — get exactly what you want from AI",
        "API Integration — connect AI models to your own apps",
        "RAG Systems — give AI your own private knowledge base",
        "Fine-tuning — train models on your specific domain",
    ]
    for i, skill in enumerate(skills):
        add_card(slide, skill,
                 Inches(0.5), Inches(3.9) + i * Inches(0.68),
                 Inches(12.33), Inches(0.55),
                 bg_color=SURFACE, text_color=WHITE, text_size=15)

    set_speaker_notes(slide, """\
SLIDE 15 — Generative AI (2 min)

"You've all used ChatGPT. Maybe Gemini. Maybe Midjourney.
But there's a massive difference between using them and building WITH them."

"At Level 5, you'll learn prompt engineering — how to get precise, professional output from AI.
API integration — connecting ChatGPT or Claude to your own applications.
RAG systems — giving AI a private knowledge base. Imagine a chatbot trained on YOUR company's data.
Fine-tuning — adapting a pre-trained model for a specific task or domain."

"This is what separates a ChatGPT user from an AI developer."
""")
    return slide


def build_slide_16(prs):
    """Level 6: AI Agents & Automation"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Level 6: AI Agents — The Future Is Autonomous")
    add_divider(slide)

    add_card(slide, "L6", Inches(0.5), Inches(1.45), Inches(0.6), Inches(0.5),
             bg_color=PURPLE, text_color=WHITE, text_size=14, bold=True)
    add_textbox(slide, "AI that doesn't just answer — it acts. End-to-end automation.",
                left=Inches(1.2), top=Inches(1.45),
                width=Inches(11.0), height=Inches(0.5),
                color=GRAY, size=15)

    steps = [
        ("Lead arrives from website form", CYAN),
        ("AI agent qualifies the enquiry", PURPLE),
        ("AI sends WhatsApp follow-up", GREEN),
        ("AI updates CRM / Google Sheets", ORANGE),
        ("AI generates performance report", CYAN),
        ("Zero human intervention needed", GREEN),
    ]
    box_w = Inches(7.5)
    box_h = Inches(0.6)
    sx = Inches(2.9)
    sy = Inches(1.55)

    for i, (text, color) in enumerate(steps):
        y = sy + i * Inches(0.82)
        add_card(slide, text, sx, y, box_w, box_h,
                 bg_color=SURFACE, text_color=color, text_size=15,
                 bold=(i == len(steps) - 1))
        if i < len(steps) - 1:
            add_textbox(slide, "↓",
                        sx + Inches(3.3), y + box_h,
                        Inches(0.9), Inches(0.22),
                        color=GRAY, size=12, align=PP_ALIGN.CENTER)

    add_card(slide, "Most in-demand AI skill in the market right now.",
             Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.45),
             bg_color=ORANGE, text_color=BG, text_size=15, bold=True)

    set_speaker_notes(slide, """\
SLIDE 16 — AI Agents (3 min)

"This is the slide I want you to remember. This is the most valuable skill in AI right now."

"An AI agent is not just a chatbot. It ACTS. It makes decisions. It connects to other tools.
It runs 24/7 without breaks or errors."

Walk through the flow:
"Lead comes in. AI agent reads it, decides if it's qualified.
If yes — WhatsApp goes out automatically. No human needed.
The CRM gets updated. The spreadsheet gets updated. A daily report is generated.
Your counsellor wakes up to warm, qualified leads in their inbox."

"Businesses are paying serious money for this.
A good AI automation consultant charges ₹50,000 to ₹2 lakh per project.
And companies are paying monthly retainers to maintain these systems."

"This is what you'll be able to build after Level 6."
""")
    return slide


def build_slide_17(prs):
    """Career Paths"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Where Can This Take You?")
    add_divider(slide)

    careers = [
        ("AI Engineer",           "₹8–25 LPA",  CYAN),
        ("ML Engineer",           "₹7–22 LPA",  PURPLE),
        ("Data Scientist",        "₹6–20 LPA",  GREEN),
        ("Data Analyst",          "₹4–12 LPA",  ORANGE),
        ("AI Consultant",         "₹8–30 LPA+", CYAN),
        ("Prompt Engineer",       "₹5–15 LPA",  PURPLE),
        ("AI Automation Spec.",   "₹6–18 LPA",  GREEN),
        ("AI Product Manager",    "₹10–30 LPA", ORANGE),
    ]
    cw, ch = Inches(2.9), Inches(1.1)
    gap_x = Inches(0.28)
    gap_y = Inches(0.3)
    sx = Inches(0.5)
    sy = Inches(1.55)

    for i, (title, salary, color) in enumerate(careers):
        row = i // 4
        col = i % 4
        x = sx + col * (cw + gap_x)
        y = sy + row * (ch + gap_y)
        add_card(slide, title, x, y, cw, Inches(0.6),
                 bg_color=SURFACE, text_color=color, text_size=14, bold=True)
        add_card(slide, salary, x, y + Inches(0.6), cw, Inches(0.45),
                 bg_color=DARK_BG2, text_color=WHITE, text_size=13)

    add_textbox(
        slide,
        "Salary ranges are indicative for India (2024–25). Entry-level to 3 years experience. Freelance/consulting can exceed these.",
        left=Inches(0.5), top=Inches(6.6),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=12, align=PP_ALIGN.CENTER, italic=True
    )

    set_speaker_notes(slide, """\
SLIDE 17 — Career Paths (3 min)

"Let's talk about where this actually takes you."

Go through each role briefly:
"AI Engineer — builds AI systems. One of the fastest-growing roles.
ML Engineer — specialises in training and deploying models.
Data Scientist — the person companies trust to make decisions from data.
Data Analyst — entry-level, but high demand. Great starting point.
AI Consultant — helps businesses implement AI. Can be freelance or corporate.
Prompt Engineer — newer role, but real and in-demand.
AI Automation Specialist — exactly what we showed in Demo 4.
AI Product Manager — leads AI product development, no coding required."

"These salary numbers are realistic — not hype. They're entry to 3-year-experience ranges.
Freelancing can take you significantly higher."

IMPORTANT: "I'm not promising you'll land these roles in 3 months.
What I'm saying is: this is where the market is going, and the earlier you build skills,
the more time compounding works in your favour."
""")
    return slide


def build_slide_18(prs):
    """Student Projects Showcase"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "What You'll Build at WorldWithWeb")
    add_divider(slide)

    projects = [
        ("Chatbot",                    "Beginner",      CYAN),
        ("Recommendation System",      "Intermediate",  ORANGE),
        ("Resume Screening AI",        "Advanced",      PURPLE),
        ("AI Interview Assistant",     "Advanced",      PURPLE),
        ("AI Sales Agent",             "Advanced",      GREEN),
        ("AI Content Generator",       "Advanced",      GREEN),
        ("AI Business Automation",     "Advanced",      ORANGE),
        ("Custom AI SaaS App",         "Capstone",      CYAN),
    ]

    pw, ph = Inches(2.9), Inches(1.05)
    gap_x = Inches(0.27)
    gap_y = Inches(0.25)
    sx = Inches(0.5)
    sy = Inches(1.55)

    for i, (name, level, color) in enumerate(projects):
        row = i // 4
        col = i % 4
        x = sx + col * (pw + gap_x)
        y = sy + row * (ph + gap_y)
        # Level badge
        add_card(slide, level, x, y, pw, Inches(0.4),
                 bg_color=SURFACE, text_color=color, text_size=11)
        # Project name
        add_card(slide, name, x, y + Inches(0.4), pw, Inches(0.6),
                 bg_color=color, text_color=BG, text_size=14, bold=True)

    add_textbox(
        slide,
        "Every project goes into your portfolio. Employers and clients see WHAT you built, not just what you studied.",
        left=Inches(0.5), top=Inches(6.55),
        width=Inches(12.33), height=Inches(0.55),
        color=GRAY, size=14, align=PP_ALIGN.CENTER, italic=True
    )

    set_speaker_notes(slide, """\
SLIDE 18 — Projects Showcase (3 min)

"Here's what you'll actually build. Not assignments. Not theory exercises. Real projects."

"A chatbot — your first AI that talks.
A recommendation system — like Netflix, built by you.
Resume screening AI — automates the first round of hiring.
AI interview assistant — helps candidates prepare for interviews.
AI sales agent — qualifies and follows up with leads automatically.
AI content generator — writes blogs, social posts, ad copy on demand.
AI business automation — connects your business tools with AI.
And a capstone SaaS app — a real, deployable product."

"When you walk into an interview and open your laptop and show these —
no one asks where you graduated from.
They ask: 'When can you start?'"
""")
    return slide


def build_slide_19(prs):
    """Why WorldWithWeb"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)
    add_title(slide, "Why WorldWithWeb?")
    add_divider(slide)

    points = [
        ("✅", "Live training with real mentors — not pre-recorded videos",          CYAN),
        ("✅", "Hands-on projects every week — build from Day 1",                    GREEN),
        ("✅", "Portfolio-first approach — show employers what you built",            CYAN),
        ("✅", "Internship guidance — real-world exposure",                          ORANGE),
        ("✅", "Interview preparation — mock interviews + resume review",             PURPLE),
        ("✅", "Placement support — referrals, job boards, industry network",         GREEN),
        ("✅", "Small batches — personal attention, not 200-person classrooms",       ORANGE),
    ]

    for i, (check, text, color) in enumerate(points):
        y = Inches(1.55) + i * Inches(0.73)
        add_card(slide, check, Inches(0.5), y, Inches(0.55), Inches(0.6),
                 bg_color=SURFACE, text_color=GREEN, text_size=18, bold=True)
        add_card(slide, text, Inches(1.15), y, Inches(11.7), Inches(0.6),
                 bg_color=SURFACE, text_color=color, text_size=15)

    set_speaker_notes(slide, """\
SLIDE 19 — Why WorldWithWeb (2 min)

"Before I tell you how to enroll, let me tell you WHY WorldWithWeb is different."

"We don't do pre-recorded videos. You get live sessions with a real mentor who can answer
your questions in real-time."

"Every week you build. No theory-only weeks. You are always making something."

"Your portfolio is built alongside the course — by the end, you have 7+ projects
you can show anyone."

"We have small batch sizes intentionally. We'd rather have 15 focused students
than 200 passive ones."

"And we're invested in your outcome — not just your enrollment. Internship guidance,
mock interviews, placement support — these aren't extras. They're part of the program."

"This isn't a YouTube playlist with a certificate at the end.
This is a structured, mentor-led, project-driven program designed to get you hired."
""")
    return slide


def build_slide_20(prs):
    """Final CTA"""
    slide = _blank_slide(prs)
    set_bg(slide, BG)

    add_textbox(
        slide,
        "The Best Time To Start AI Was Yesterday.",
        left=Inches(0.5), top=Inches(0.7),
        width=Inches(12.33), height=Inches(1.1),
        color=CYAN, size=40, bold=True, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide,
        "The Second Best Time Is Today.",
        left=Inches(0.5), top=Inches(1.75),
        width=Inches(12.33), height=Inches(0.8),
        color=WHITE, size=30, bold=False, align=PP_ALIGN.CENTER
    )

    # Divider line
    shape = slide.shapes.add_shape(1, Inches(3.0), Inches(2.7), Inches(7.33), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()

    add_textbox(
        slide,
        "One year from now — will you be using AI,  or building it?",
        left=Inches(0.5), top=Inches(2.85),
        width=Inches(12.33), height=Inches(0.65),
        color=ORANGE, size=20, bold=True, align=PP_ALIGN.CENTER
    )

    # Enrollment details box
    add_card(slide,
             "Next Batch: [DATE]     |     Fee: [AMOUNT]     |     Seats: Limited",
             Inches(1.5), Inches(3.7), Inches(10.33), Inches(0.75),
             bg_color=SURFACE, text_color=WHITE, text_size=18)

    add_textbox(
        slide, "Enroll Now  →  Talk to us after this session",
        left=Inches(0.5), top=Inches(4.6),
        width=Inches(12.33), height=Inches(0.55),
        color=GREEN, size=18, bold=True, align=PP_ALIGN.CENTER
    )

    # Contact / social
    add_textbox(
        slide, "📱  WhatsApp: [NUMBER]     🌐  worldwithweb.in     📧  [EMAIL]",
        left=Inches(0.5), top=Inches(5.3),
        width=Inches(12.33), height=Inches(0.5),
        color=GRAY, size=15, align=PP_ALIGN.CENTER
    )

    add_textbox(
        slide, "WorldWithWeb — Learn Tech. Build Real Things.",
        left=Inches(0.5), top=Inches(6.75),
        width=Inches(12.33), height=Inches(0.45),
        color=ORANGE, size=14, bold=True, align=PP_ALIGN.CENTER
    )

    set_speaker_notes(slide, """\
SLIDE 20 — Final CTA (5 min)

"I want to leave you with one thought."

[Read the headline slowly]
"The best time to start AI was yesterday. The second best time is today."

[Pause]

"In this room right now, you have a choice.
One year from now, you can be someone who uses AI tools made by others —
OR you can be someone who BUILT those tools, who gets PAID to build them, who is in demand."

"The difference between those two people starts with one decision."

[Point to enrollment box]
"Our next batch starts on [DATE]. We have limited seats — we keep it small intentionally.
The fee is [AMOUNT]."

"I'm going to be right here after this session.
If you want to talk about the course, your goals, what track is right for you — come find me.
No pressure. No script. Just a conversation."

[Closing line]
"Thank you. Let's build something."

POST-SESSION:
- Collect names and WhatsApp numbers
- Note which demos excited each student most
- Share WorldWithWeb contact and batch details
- Follow up within 24 hours
""")
    return slide


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    builders = [
        build_slide_01, build_slide_02, build_slide_03, build_slide_04,
        build_slide_05, build_slide_06, build_slide_07, build_slide_08,
        build_slide_09, build_slide_10, build_slide_11, build_slide_12,
        build_slide_13, build_slide_14, build_slide_15, build_slide_16,
        build_slide_17, build_slide_18, build_slide_19, build_slide_20,
    ]

    for i, fn in enumerate(builders, 1):
        print(f"  Building slide {i:02d}/20 — {fn.__doc__}")
        fn(prs)

    out_dir = os.path.join(os.path.dirname(__file__), "..", "presentations")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "WorldWithWeb_AI_Motivation_Session.pptx")
    prs.save(out_path)
    print(f"\n✅  Saved → {os.path.abspath(out_path)}")


if __name__ == "__main__":
    main()
