"""Generate the WorldWithWeb AI Trends Session PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors (same as Demo Session PPT)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x1C)
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)
ACCENT_RED = RGBColor(0xFF, 0x45, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_CARD = RGBColor(0x25, 0x25, 0x40)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


# =====================================================
# SLIDE 1 - Title
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.8), Inches(4.3), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

tb = add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.5))
set_text(tb.text_frame, "WorldWithWeb", size=56, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

tb2 = add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.5))
set_text(tb2.text_frame, "AI Trends & The Future of Skills", size=36, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

tb3 = add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1))
set_text(tb3.text_frame, "How AI is changing every career  |  Session 2", size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================
# SLIDE 2 - AI Is Already in Your Life
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

tb = add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(2))
set_text(tb.text_frame, "AI Is Already in Your Life", size=48, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

tb2 = add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.5))
set_text(tb2.text_frame, '"Name one thing you used today that has AI in it."', size=32, color=WHITE, alignment=PP_ALIGN.CENTER)

# Example cards
examples = [
    ("Google Maps", ACCENT_BLUE),
    ("Instagram", ACCENT_PURPLE),
    ("YouTube", ACCENT_RED),
    ("Autocomplete", ACCENT_ORANGE),
    ("Spotify", ACCENT_GREEN),
]
for i, (name, color) in enumerate(examples):
    x = Inches(1.0) + Inches(2.3) * i
    card = add_shape(slide, x, Inches(5.2), Inches(2.0), Inches(1.0), DARK_CARD, color)
    tb_c = add_text_box(slide, x, Inches(5.35), Inches(2.0), Inches(0.8))
    set_text(tb_c.text_frame, name, size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)


# =====================================================
# SLIDE 3 - The AI Explosion: 2020-2026
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Accent bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), H)
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_PURPLE
bar.line.fill.background()

tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(1))
set_text(tb.text_frame, "The AI Explosion: 2020-2026", size=40, bold=True, color=WHITE)

timeline = [
    ("2020", "GPT-3", "AI could write like a human", ACCENT_BLUE),
    ("2021", "Copilot + DALL-E", "AI started coding & creating images", ACCENT_GREEN),
    ("2022", "ChatGPT", "AI became mainstream overnight", ACCENT_ORANGE),
    ("2023", "GPT-4, Gemini", "Multimodal AI -- text, images, code", ACCENT_PURPLE),
    ("2024", "Sora, AI Agents", "AI creates videos, works autonomously", ACCENT_RED),
    ("2025-26", "AI Everywhere", "Not a tool anymore -- it's infrastructure", ACCENT_GREEN),
]

for i, (year, name, desc, color) in enumerate(timeline):
    y = Inches(1.6) + Inches(0.95) * i
    card = add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.78), DARK_CARD, color)
    tb_y = add_text_box(slide, Inches(1.1), y + Emu(Inches(0.12).emu), Inches(1.8), Inches(0.6))
    set_text(tb_y.text_frame, year, size=20, bold=True, color=color)
    tb_n = add_text_box(slide, Inches(3.0), y + Emu(Inches(0.12).emu), Inches(3.5), Inches(0.6))
    set_text(tb_n.text_frame, name, size=20, bold=True, color=WHITE)
    tb_d = add_text_box(slide, Inches(6.8), y + Emu(Inches(0.12).emu), Inches(5.2), Inches(0.6))
    set_text(tb_d.text_frame, desc, size=18, color=LIGHT_GRAY)


# =====================================================
# SLIDE 4 - AI in Numbers
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

tb = add_text_box(slide, Inches(1.5), Inches(0.5), Inches(10.3), Inches(1))
set_text(tb.text_frame, "AI in Numbers", size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

stats = [
    ("80%", "of companies will\nuse AI by 2026", ACCENT_BLUE),
    ("97M", "new AI-related\njobs globally", ACCENT_GREEN),
    ("$17B", "India's AI market\nby 2027", ACCENT_ORANGE),
]

for i, (num, desc, color) in enumerate(stats):
    x = Inches(0.8) + Inches(4.1) * i
    card = add_shape(slide, x, Inches(2.5), Inches(3.7), Inches(3.5), DARK_CARD, color)
    tb_n = add_text_box(slide, x, Inches(2.9), Inches(3.7), Inches(1.5))
    set_text(tb_n.text_frame, num, size=56, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    tb_d = add_text_box(slide, x, Inches(4.5), Inches(3.7), Inches(1.5))
    set_text(tb_d.text_frame, desc, size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

tb_q = add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.8))
set_text(tb_q.text_frame, '"This is not hype. This is infrastructure. Like the internet was in 2005."', size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================
# SLIDE 5 - AI in Coding & Development
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_GREEN
bar.line.fill.background()

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6))
set_text(tb.text_frame, "AI ACROSS FIELDS  |  1 of 5", size=16, color=ACCENT_GREEN, bold=True)

tb2 = add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(1.2))
set_text(tb2.text_frame, "AI in Coding & App Development", size=40, bold=True, color=WHITE)

# Bullets
tb3 = add_text_box(slide, Inches(0.8), Inches(2.5), Inches(6.5), Inches(3.5))
tf = tb3.text_frame
tf.word_wrap = True
set_text(tf, "GitHub Copilot -- AI writes code alongside you", size=22, color=WHITE)
add_paragraph(tf, "Cursor, Windsurf -- AI-powered code editors", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "AI Code Review -- catches bugs before they ship", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "Automated Testing -- AI writes test cases", size=22, color=WHITE, space_before=Pt(12))

# Right card
card = add_shape(slide, Inches(8), Inches(2.5), Inches(4.5), Inches(3), DARK_CARD, ACCENT_GREEN)
tb_c = add_text_box(slide, Inches(8.3), Inches(2.8), Inches(4), Inches(2.5))
tf2 = tb_c.text_frame
set_text(tf2, "The Reality", size=20, color=ACCENT_GREEN, bold=True)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "AI writes code.", size=22, color=WHITE)
add_paragraph(tf2, "You decide WHAT", size=22, color=WHITE)
add_paragraph(tf2, "to build.", size=22, color=WHITE)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "That's the Builder.", size=20, color=ACCENT_GREEN, bold=True)

# Course link
tb_link = add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6))
set_text(tb_link.text_frame, "WorldWithWeb Track: Python & AI Program", size=18, color=ACCENT_GREEN, bold=True)


# =====================================================
# SLIDE 6 - AI in Digital Marketing
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_ORANGE
bar.line.fill.background()

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6))
set_text(tb.text_frame, "AI ACROSS FIELDS  |  2 of 5", size=16, color=ACCENT_ORANGE, bold=True)

tb2 = add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(1.2))
set_text(tb2.text_frame, "AI in Digital Marketing & SEO", size=40, bold=True, color=WHITE)

tb3 = add_text_box(slide, Inches(0.8), Inches(2.5), Inches(6.5), Inches(3.5))
tf = tb3.text_frame
tf.word_wrap = True
set_text(tf, "ChatGPT for copywriting -- ads, blogs, emails in seconds", size=22, color=WHITE)
add_paragraph(tf, "AI Ad Targeting -- Facebook, Google find your customers", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "SEO Tools (Surfer, Jasper) -- optimize for Google rankings", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "AI Chatbots -- 24/7 customer support, zero hiring", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "Social Media Analytics -- AI predicts viral content", size=22, color=WHITE, space_before=Pt(12))

# Right card
card = add_shape(slide, Inches(8), Inches(2.5), Inches(4.5), Inches(3), DARK_CARD, ACCENT_ORANGE)
tb_c = add_text_box(slide, Inches(8.3), Inches(2.8), Inches(4), Inches(2.5))
tf2 = tb_c.text_frame
set_text(tf2, "The Reality", size=20, color=ACCENT_ORANGE, bold=True)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "The marketer who", size=22, color=WHITE)
add_paragraph(tf2, "uses AI does 10x.", size=22, color=WHITE)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "The one who doesn't", size=20, color=LIGHT_GRAY)
add_paragraph(tf2, "gets left behind.", size=20, color=LIGHT_GRAY)

tb_link = add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6))
set_text(tb_link.text_frame, "WorldWithWeb Track: Digital Marketing Program", size=18, color=ACCENT_ORANGE, bold=True)


# =====================================================
# SLIDE 7 - AI in Content & Design
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_PURPLE
bar.line.fill.background()

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6))
set_text(tb.text_frame, "AI ACROSS FIELDS  |  3 of 5", size=16, color=ACCENT_PURPLE, bold=True)

tb2 = add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(1.2))
set_text(tb2.text_frame, "AI in Content Creation & Design", size=40, bold=True, color=WHITE)

tb3 = add_text_box(slide, Inches(0.8), Inches(2.5), Inches(6.5), Inches(3.5))
tf = tb3.text_frame
tf.word_wrap = True
set_text(tf, "Midjourney, DALL-E -- stunning images from text prompts", size=22, color=WHITE)
add_paragraph(tf, "Canva AI -- auto-generate designs, remove backgrounds", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "CapCut, Runway -- AI video editing, auto-captions", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "AI Reels Scripts -- content ideas and scripts instantly", size=22, color=WHITE, space_before=Pt(12))

card = add_shape(slide, Inches(8), Inches(2.5), Inches(4.5), Inches(3), DARK_CARD, ACCENT_PURPLE)
tb_c = add_text_box(slide, Inches(8.3), Inches(2.8), Inches(4), Inches(2.5))
tf2 = tb_c.text_frame
set_text(tf2, "The Reality", size=20, color=ACCENT_PURPLE, bold=True)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "Create more.", size=24, color=WHITE)
add_paragraph(tf2, "Create better.", size=24, color=WHITE)
add_paragraph(tf2, "Create faster.", size=24, color=WHITE)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "AI is the multiplier.", size=20, color=ACCENT_PURPLE, bold=True)

tb_link = add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6))
set_text(tb_link.text_frame, "WorldWithWeb Track: Content & Design Bootcamp", size=18, color=ACCENT_PURPLE, bold=True)


# =====================================================
# SLIDE 8 - AI in Cybersecurity
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_RED
bar.line.fill.background()

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6))
set_text(tb.text_frame, "AI ACROSS FIELDS  |  4 of 5", size=16, color=ACCENT_RED, bold=True)

tb2 = add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(1.2))
set_text(tb2.text_frame, "AI in Cybersecurity", size=40, bold=True, color=WHITE)

tb3 = add_text_box(slide, Inches(0.8), Inches(2.5), Inches(6.5), Inches(3.5))
tf = tb3.text_frame
tf.word_wrap = True
set_text(tf, "AI Threat Detection -- catches attacks humans would miss", size=22, color=WHITE)
add_paragraph(tf, "Anomaly Detection -- spots unusual network patterns", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "Automated Pen Testing -- finds vulnerabilities first", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "The Arms Race -- AI vs AI in cyber warfare", size=22, color=WHITE, space_before=Pt(12))

card = add_shape(slide, Inches(8), Inches(2.5), Inches(4.5), Inches(3), DARK_CARD, ACCENT_RED)
tb_c = add_text_box(slide, Inches(8.3), Inches(2.8), Inches(4), Inches(2.5))
tf2 = tb_c.text_frame
set_text(tf2, "The Reality", size=20, color=ACCENT_RED, bold=True)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "Security teams can't", size=22, color=WHITE)
add_paragraph(tf2, "survive without AI.", size=22, color=WHITE)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "Highest-paying tech", size=20, color=LIGHT_GRAY)
add_paragraph(tf2, "careers start here.", size=20, color=LIGHT_GRAY)

tb_link = add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6))
set_text(tb_link.text_frame, "WorldWithWeb: Python Cybersecurity Projects", size=18, color=ACCENT_RED, bold=True)


# =====================================================
# SLIDE 9 - AI in Education
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

tb = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6))
set_text(tb.text_frame, "AI ACROSS FIELDS  |  5 of 5", size=16, color=ACCENT_BLUE, bold=True)

tb2 = add_text_box(slide, Inches(0.8), Inches(1), Inches(11), Inches(1.2))
set_text(tb2.text_frame, "AI in Education & Learning", size=40, bold=True, color=WHITE)

tb3 = add_text_box(slide, Inches(0.8), Inches(2.5), Inches(6.5), Inches(3.5))
tf = tb3.text_frame
tf.word_wrap = True
set_text(tf, "Personalized Learning -- AI adapts to each student's level", size=22, color=WHITE)
add_paragraph(tf, "AI Tutors -- explanations in any style, any time", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "Smart Assessments -- AI generates and grades quizzes", size=22, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "Language Learning -- Duolingo's AI adapts in real-time", size=22, color=WHITE, space_before=Pt(12))

card = add_shape(slide, Inches(8), Inches(2.5), Inches(4.5), Inches(3), DARK_CARD, ACCENT_BLUE)
tb_c = add_text_box(slide, Inches(8.3), Inches(2.8), Inches(4), Inches(2.5))
tf2 = tb_c.text_frame
set_text(tf2, "The Reality", size=20, color=ACCENT_BLUE, bold=True)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "You're going to build", size=22, color=WHITE)
add_paragraph(tf2, "exactly this.", size=22, color=WHITE)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "Quiz Master. AI Tutor.", size=20, color=ACCENT_BLUE, bold=True)
add_paragraph(tf2, "Not just projects --", size=18, color=LIGHT_GRAY)
add_paragraph(tf2, "the future of education.", size=18, color=LIGHT_GRAY)

tb_link = add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6))
set_text(tb_link.text_frame, "WorldWithWeb Projects: Quiz Master, AI Tutor", size=18, color=ACCENT_BLUE, bold=True)


# =====================================================
# SLIDE 10 - 3 Types of AI Careers
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

tb = add_text_box(slide, Inches(1.5), Inches(0.3), Inches(10.3), Inches(1))
set_text(tb.text_frame, "3 Types of AI Careers", size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# AI User
card1 = add_shape(slide, Inches(0.8), Inches(1.6), Inches(3.7), Inches(4.5), DARK_CARD, MID_GRAY)
tb1 = add_text_box(slide, Inches(1.1), Inches(1.9), Inches(3.2), Inches(4))
tf = tb1.text_frame
set_text(tf, "Type 1", size=18, color=MID_GRAY, bold=True)
add_paragraph(tf, "AI USER", size=32, bold=True, color=MID_GRAY)
add_paragraph(tf, "", size=8)
add_paragraph(tf, "Uses AI tools daily", size=18, color=LIGHT_GRAY)
add_paragraph(tf, "Marketers, designers,", size=18, color=LIGHT_GRAY)
add_paragraph(tf, "writers, students", size=18, color=LIGHT_GRAY)
add_paragraph(tf, "", size=8)
add_paragraph(tf, "The minimum bar.", size=18, color=MID_GRAY)
add_paragraph(tf, "Everyone will be here.", size=18, color=MID_GRAY)

# AI Builder
card2 = add_shape(slide, Inches(4.8), Inches(1.6), Inches(3.7), Inches(4.5), DARK_CARD, ACCENT_BLUE)
tb2 = add_text_box(slide, Inches(5.1), Inches(1.9), Inches(3.2), Inches(4))
tf2 = tb2.text_frame
set_text(tf2, "Type 2", size=18, color=ACCENT_BLUE, bold=True)
add_paragraph(tf2, "AI BUILDER", size=32, bold=True, color=ACCENT_BLUE)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "Builds AI-powered products", size=18, color=WHITE)
add_paragraph(tf2, "Developers, data scientists,", size=18, color=WHITE)
add_paragraph(tf2, "automation engineers", size=18, color=WHITE)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "Where the money is.", size=18, color=ACCENT_BLUE)
add_paragraph(tf2, "Where the demand is.", size=18, color=ACCENT_BLUE)

# AI Strategist
card3 = add_shape(slide, Inches(8.8), Inches(1.6), Inches(3.7), Inches(4.5), DARK_CARD, ACCENT_GREEN)
tb3 = add_text_box(slide, Inches(9.1), Inches(1.9), Inches(3.2), Inches(4))
tf3 = tb3.text_frame
set_text(tf3, "Type 3", size=18, color=ACCENT_GREEN, bold=True)
add_paragraph(tf3, "AI STRATEGIST", size=30, bold=True, color=ACCENT_GREEN)
add_paragraph(tf3, "", size=8)
add_paragraph(tf3, "Decides how organizations", size=18, color=WHITE)
add_paragraph(tf3, "adopt AI", size=18, color=WHITE)
add_paragraph(tf3, "PMs, CTOs, founders", size=18, color=WHITE)
add_paragraph(tf3, "", size=8)
add_paragraph(tf3, "Where the leadership is.", size=18, color=ACCENT_GREEN)
add_paragraph(tf3, "Where the vision is.", size=18, color=ACCENT_GREEN)

tb_q = add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.8))
set_text(tb_q.text_frame, "WorldWithWeb moves you from User to Builder. In any field.", size=22, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# =====================================================
# SLIDE 11 - Where Does YOUR Course Fit?
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

tb = add_text_box(slide, Inches(1.5), Inches(0.3), Inches(10.3), Inches(1))
set_text(tb.text_frame, "Where Does YOUR Course Fit?", size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Python Track card
card1 = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.5), DARK_CARD, ACCENT_BLUE)
tb1 = add_text_box(slide, Inches(1.2), Inches(2.1), Inches(5), Inches(4))
tf = tb1.text_frame
set_text(tf, "Python & AI Track", size=28, bold=True, color=ACCENT_BLUE)
add_paragraph(tf, "", size=8)
add_paragraph(tf, "Build real AI-powered applications", size=20, color=WHITE)
add_paragraph(tf, "7 projects in 6 weeks", size=20, color=WHITE)
add_paragraph(tf, "Chatbots, games, security tools", size=20, color=WHITE)
add_paragraph(tf, "Learn to work WITH AI as a co-builder", size=20, color=WHITE)
add_paragraph(tf, "", size=10)
add_paragraph(tf, "Career paths:", size=18, color=MID_GRAY)
add_paragraph(tf, "Developer, AI Engineer, Cybersecurity,", size=18, color=LIGHT_GRAY)
add_paragraph(tf, "Data Science, Automation", size=18, color=LIGHT_GRAY)

# Digital Marketing Track card
card2 = add_shape(slide, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.5), DARK_CARD, ACCENT_ORANGE)
tb2 = add_text_box(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(4))
tf2 = tb2.text_frame
set_text(tf2, "Digital Marketing Track", size=28, bold=True, color=ACCENT_ORANGE)
add_paragraph(tf2, "", size=8)
add_paragraph(tf2, "Master AI-powered marketing tools", size=20, color=WHITE)
add_paragraph(tf2, "SEO, social media, content creation", size=20, color=WHITE)
add_paragraph(tf2, "AI chatbots, ad targeting, analytics", size=20, color=WHITE)
add_paragraph(tf2, "Build campaigns that outperform", size=20, color=WHITE)
add_paragraph(tf2, "", size=10)
add_paragraph(tf2, "Career paths:", size=18, color=MID_GRAY)
add_paragraph(tf2, "Digital Marketer, Content Strategist,", size=18, color=LIGHT_GRAY)
add_paragraph(tf2, "SEO Specialist, Social Media Manager", size=18, color=LIGHT_GRAY)

# Bottom message
tb_b = add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.7))
set_text(tb_b.text_frame, "Both tracks make you an AI Builder, not just an AI User.", size=24, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# =====================================================
# SLIDE 12 - What WorldWithWeb Students Build
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

tb = add_text_box(slide, Inches(1.5), Inches(0.3), Inches(10.3), Inches(1))
set_text(tb.text_frame, "What WorldWithWeb Students Build", size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

tb_sub = add_text_box(slide, Inches(1.5), Inches(1.1), Inches(10.3), Inches(0.6))
set_text(tb_sub.text_frame, "Real projects. Real portfolio. Real skills.", size=24, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

projects = [
    ("Quiz Master", "AI-powered quiz on any topic", ACCENT_BLUE),
    ("Password Fortress", "Security analyzer & cracker", ACCENT_ORANGE),
    ("AI Chatbot", "Custom personalities & models", ACCENT_PURPLE),
    ("Network Sentinel", "Cybersecurity scanner", ACCENT_RED),
    ("Space Defender", "Playable space shooter game", ACCENT_GREEN),
    ("AI Story Forge", "AI-generated choose-your-adventure", ACCENT_BLUE),
]

for i, (name, desc, color) in enumerate(projects):
    row = i // 3
    col = i % 3
    x = Inches(0.8) + Inches(4.1) * col
    y = Inches(2.2) + Inches(2.3) * row
    card = add_shape(slide, x, y, Inches(3.7), Inches(1.8), DARK_CARD, color)
    tb_n = add_text_box(slide, x + Emu(Inches(0.3).emu), y + Emu(Inches(0.3).emu), Inches(3.2), Inches(0.8))
    set_text(tb_n.text_frame, name, size=24, bold=True, color=color)
    tb_d = add_text_box(slide, x + Emu(Inches(0.3).emu), y + Emu(Inches(1.0).emu), Inches(3.2), Inches(0.7))
    set_text(tb_d.text_frame, desc, size=18, color=LIGHT_GRAY)


# =====================================================
# SLIDE 13 - Which AI Trend Excites You?
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

tb = add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(2))
set_text(tb.text_frame, "Which AI Trend\nExcites You Most?", size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Field options in their colors
fields = [
    ("Coding?", ACCENT_GREEN),
    ("Marketing?", ACCENT_ORANGE),
    ("Design?", ACCENT_PURPLE),
    ("Security?", ACCENT_RED),
    ("Education?", ACCENT_BLUE),
]
for i, (field, color) in enumerate(fields):
    x = Inches(1.0) + Inches(2.3) * i
    card = add_shape(slide, x, Inches(4.5), Inches(2.0), Inches(1.0), DARK_CARD, color)
    tb_f = add_text_box(slide, x, Inches(4.65), Inches(2.0), Inches(0.7))
    set_text(tb_f.text_frame, field, size=22, bold=True, color=color, alignment=PP_ALIGN.CENTER)

tb3 = add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.8))
set_text(tb3.text_frame, "[ Open discussion -- each person picks a field ]", size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# =====================================================
# SLIDE 14 - Closing
# =====================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.5), Inches(4.3), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

tb = add_text_box(slide, Inches(1.5), Inches(1), Inches(10.3), Inches(1.5))
set_text(tb.text_frame, "WorldWithWeb", size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

tb2 = add_text_box(slide, Inches(1.5), Inches(3), Inches(10.3), Inches(1))
set_text(tb2.text_frame, "AI is not the future. It's the now.", size=30, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

tb3 = add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.5))
tf = tb3.text_frame
set_text(tf, "worldwithweb.com", size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "", size=10)
add_paragraph(tf, "The only question is: will you be ready?", size=28, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# Save
output_path = r"c:\Users\naaz.verma\personal\python_projects\WorldWithWeb_AI_Trends_Session.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
