"""
Generate an academic-style PowerPoint presentation for:
"Learning an Interpretable Traffic Signal Control Policy"
by James Ault, Josiah P. Hanna, Guni Sharon (arXiv: 1912.11023v2)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Academic colour palette ──
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_WHITE   = RGBColor(0xFB, 0xFB, 0xFB)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
NAVY       = RGBColor(0x1B, 0x3A, 0x5C)   # primary heading color
BLUE       = RGBColor(0x2B, 0x6C, 0xB3)   # accent / links
DARK_BLUE  = RGBColor(0x15, 0x2D, 0x4A)   # title bar
RED_ACC    = RGBColor(0xC0, 0x39, 0x2B)    # for emphasis / limitations
GREEN_ACC  = RGBColor(0x27, 0x7A, 0x3E)    # for positive results
BOX_BG     = RGBColor(0xEE, 0xF2, 0xF7)   # light blue-gray box fill
BOX_BG2    = RGBColor(0xFD, 0xF0, 0xE0)   # light warm box fill
BOX_BORDER = RGBColor(0xB0, 0xC4, 0xDE)   # light steel blue


def set_slide_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_GRAY, spacing=Pt(8), font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return tf


def add_bar(slide, left, top, width, height, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_box(slide, left, top, width, height, text, font_size=14,
            fill_color=BOX_BG, text_color=DARK_GRAY, border_color=BOX_BORDER,
            alignment=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "Calibri"
    return shape


def slide_header(slide, title, subtitle=None):
    """Add standard academic slide header with navy bar and title."""
    set_slide_bg(slide)
    add_bar(slide, 0, 0, 13.333, 0.9, DARK_BLUE)
    add_textbox(slide, 0.6, 0.15, 12, 0.6, title,
                font_size=28, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, 0.6, 0.55, 12, 0.35, subtitle,
                    font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD))
    # thin accent line below header
    add_bar(slide, 0, 0.9, 13.333, 0.04, BLUE)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ═══════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_bar(s, 0, 0, 13.333, 2.8, DARK_BLUE)

add_textbox(s, 1.5, 0.6, 10.3, 1.2,
            "Learning an Interpretable\nTraffic Signal Control Policy",
            font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bar(s, 5.5, 2.05, 2.3, 0.03, RGBColor(0xBB, 0xCC, 0xDD))

add_textbox(s, 1.5, 2.2, 10.3, 0.5,
            "James Ault  |  Josiah P. Hanna  |  Guni Sharon",
            font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 3.4, 10.3, 0.5,
            "arXiv: 1912.11023v2  |  2020",
            font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 4.5, 10.3, 0.5,
            "Presented by: Naaz Verma",
            font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 5.2, 10.3, 0.5,
            "Assignment Presentation",
            font_size=15, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Footer line
add_bar(s, 0, 7.2, 13.333, 0.03, BLUE)

# ═══════════════════════════════════════════════
# SLIDE 2 — Outline
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Presentation Outline")

agenda = [
    "1.   Motivation & Objective of the Paper",
    "2.   Key Techniques & Algorithmic Approaches",
    "3.   Experimental Setup",
    "4.   Main Results & Interpretation",
    "5.   Limitations & Future Scope",
    "6.   Conclusion & Key Takeaways",
]
add_bullet_list(s, 1.2, 1.5, 8, 5, agenda, font_size=20, color=DARK_GRAY, spacing=Pt(18))

# ═══════════════════════════════════════════════
# SLIDE 3 — Motivation & Problem
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Motivation & Problem Statement")

add_textbox(s, 0.6, 1.2, 6.2, 0.5, "The Interpretability Gap in Traffic Signal Control",
            font_size=18, color=NAVY, bold=True)

items = [
    "\u2022  Traffic congestion costs billions annually in fuel, time, and emissions",
    "\u2022  Reinforcement Learning with DNNs can reduce vehicle delay by up to 73%",
    "\u2022  However, DNNs are black boxes \u2014 decisions cannot be explained",
    "\u2022  Government agencies require interpretable, auditable controllers",
    "\u2022  Current RL approaches are impractical for real-world deployment",
]
add_bullet_list(s, 0.8, 1.8, 6.2, 3.5, items, font_size=15, color=DARK_GRAY, spacing=Pt(10))

add_box(s, 7.4, 1.2, 5.3, 2.3,
        "Core Research Question\n\n"
        "Can we design traffic signal controllers\n"
        "that are BOTH high-performing AND\n"
        "human-interpretable?\n\n"
        "Performance (DNN)  vs.  Interpretability",
        font_size=15, fill_color=BOX_BG, text_color=DARK_GRAY,
        border_color=BLUE, alignment=PP_ALIGN.CENTER)

add_box(s, 7.4, 3.8, 5.3, 3.0,
        "Why Interpretability is Non-Negotiable\n\n"
        "\u2022  Liability: Who is responsible if an accident occurs?\n"
        "\u2022  Regulation: Agencies must approve control logic\n"
        "\u2022  Trust: Engineers need to understand & tune the system\n"
        "\u2022  Safety: Black-box failures can be catastrophic\n\n"
        "Traffic signals are safety-critical infrastructure.",
        font_size=14, fill_color=BOX_BG2, text_color=DARK_GRAY,
        border_color=RGBColor(0xD4, 0xA0, 0x60))

# ═══════════════════════════════════════════════
# SLIDE 4 — Objective & Contributions
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Paper Objective & Contributions")

add_box(s, 0.6, 1.3, 12.1, 1.0,
        "Objective: Design interpretable, regulatable control policies for traffic signals that match "
        "deep neural network performance while remaining fully human-understandable and tunable.",
        font_size=16, fill_color=BOX_BG, text_color=NAVY,
        border_color=NAVY, alignment=PP_ALIGN.CENTER)

add_textbox(s, 0.6, 2.6, 5, 0.5, "Six Key Contributions:",
            font_size=18, color=NAVY, bold=True)

contributions = [
    "1.  Formally define a \"regulatable\" control function for signal control",
    "2.  Compare regulatable functions against DNN-based policies",
    "3.  Study three optimization methods: CMA-ES, PPO, Deep Q-Learning",
    "4.  Develop three novel DQN variants using regulatable functions",
    "      \u2192 DRQ, DRSQ, DRHQ (core technical contribution)",
    "5.  Evaluate on simulations of real intersections with observed traffic demand",
    "6.  Compare against deployed actuated controllers (practical baseline)",
]
add_bullet_list(s, 0.8, 3.2, 11.5, 4, contributions, font_size=15, color=DARK_GRAY, spacing=Pt(9))

# ═══════════════════════════════════════════════
# SLIDE 5 — Regulatable Control Function
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Key Concept: Regulatable Control Function")

add_box(s, 0.6, 1.3, 5.8, 2.2,
        "Formal Definition\n\n"
        "A precedence function g(s, \u03a6; \u03b8) is regulatable\n"
        "if for every state variable s[i]:\n\n"
        "    \u2202g/\u2202s[i]  \u2265  0   for all states,   OR\n"
        "    \u2202g/\u2202s[i]  \u2264  0   for all states\n\n"
        "i.e., monotonic input-output relationships.",
        font_size=14, fill_color=BOX_BG, text_color=DARK_GRAY, border_color=NAVY)

add_box(s, 6.8, 1.3, 5.9, 2.2,
        "Practical Intuition\n\n"
        "\"Green was given to Phase 4 BECAUSE\n"
        "stopped Southbound vehicles increased\n"
        "while Eastbound queue decreased.\"\n\n"
        "Every decision has a human-readable reason.\n"
        "Engineers can adjust individual weights.",
        font_size=14, fill_color=BOX_BG2, text_color=DARK_GRAY,
        border_color=RGBColor(0xD4, 0xA0, 0x60))

add_textbox(s, 0.6, 3.9, 5.8, 0.4, "State Variables (per traffic phase):",
            font_size=16, color=NAVY, bold=True)

state_vars = [
    "1.  Stopped vehicles count",
    "2.  Approaching vehicles count",
    "3.  Cumulative stopped time",
    "4.  Average stopped time",
    "5.  Average queue length (stopped vehicles / lane count)",
    "6.  Average approaching vehicle speed",
]
add_bullet_list(s, 0.8, 4.4, 5.5, 2.5, state_vars, font_size=14, color=DARK_GRAY, spacing=Pt(5))

add_textbox(s, 6.8, 3.9, 5.9, 0.4, "Polynomial Precedence Function:",
            font_size=16, color=NAVY, bold=True)

add_box(s, 6.8, 4.4, 5.9, 1.0,
        "g(s, \u03a6; \u03b8') = \u03a3\u03c6\u2208\u03a6 \u03a3i=1..6  (wi \u00b7 s[i])^pi   \u00d7   \u03a3j=1..4  (w'j \u00b7 fj)^p'j",
        font_size=14, fill_color=BOX_BG, text_color=NAVY,
        border_color=BLUE, alignment=PP_ALIGN.CENTER)

props = [
    "\u2022  Weights (w) and exponents (p) per state variable per phase",
    "\u2022  For 8-phase intersection: 256 total tunable parameters",
    "\u2022  Proven regulatable via Lemma 1 (monotonic partial derivatives)",
]
add_bullet_list(s, 7.0, 5.6, 5.5, 1.5, props, font_size=13, color=MID_GRAY, spacing=Pt(5))

# ═══════════════════════════════════════════════
# SLIDE 6 — Three DQN Variants
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Algorithmic Approach: Three Deep Regulatable Q-Learning Variants")

add_box(s, 0.6, 1.2, 12.1, 0.7,
        "Core Idea:  Train a powerful DQN (black-box) first, then train the interpretable function G to imitate it.",
        font_size=16, fill_color=BOX_BG, text_color=NAVY,
        border_color=NAVY, alignment=PP_ALIGN.CENTER)

# DRQ
add_box(s, 0.6, 2.2, 3.8, 3.6,
        "DRQ\nDeep Regulatable Q-Learning\n\n"
        "Goal: G(s,a) = Q(s,a)\n"
        "(Match exact Q-values)\n\n"
        "Loss: Squared error\n"
        "(y - G(s,a;\u03b8))\u00b2\n\n"
        "Problem: Polynomial cannot\nreplicate full DNN capacity\n\n"
        "Result: Worst performer",
        font_size=13, fill_color=RGBColor(0xFC, 0xEB, 0xEB),
        text_color=DARK_GRAY, border_color=RED_ACC,
        alignment=PP_ALIGN.CENTER)

# DRSQ
add_box(s, 4.8, 2.2, 3.8, 3.6,
        "DRSQ\nDeep Regulatable Softmax Q\n\n"
        "Goal: G(s,\u00b7) \u221d Q(s,\u00b7)\n"
        "(Match relative ranking)\n\n"
        "Loss: Cross-entropy between\nsoftmax(Q) and softmax(G)\n\n"
        "Insight: Proportional\nequivalence is sufficient\n\n"
        "Result: Good performer",
        font_size=13, fill_color=BOX_BG2,
        text_color=DARK_GRAY, border_color=RGBColor(0xD4, 0xA0, 0x60),
        alignment=PP_ALIGN.CENTER)

# DRHQ
add_box(s, 9.0, 2.2, 3.8, 3.6,
        "DRHQ  (Best)\nDeep Regulatable Hardmax Q\n\n"
        "Goal: argmax G = argmax Q\n"
        "(Match only the winning action)\n\n"
        "Loss: Cross-entropy between\none-hot(argmax Q) & softmax(G)\n\n"
        "Insight: Maximum flexibility,\nonly policy equivalence needed\n\n"
        "Result: Best performer",
        font_size=13, fill_color=RGBColor(0xE8, 0xF5, 0xEB),
        text_color=DARK_GRAY, border_color=GREEN_ACC,
        alignment=PP_ALIGN.CENTER)

# Progression arrow
add_box(s, 0.6, 6.1, 12.1, 0.6,
        "Progression:   Match exact values (hard)   \u2192   Match rankings (easier)   \u2192   Match winner only (easiest & best)",
        font_size=15, fill_color=BOX_BG, text_color=NAVY,
        border_color=BLUE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 7 — CMA-ES & PPO
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Alternative Optimization Methods: CMA-ES & PPO")

# CMA-ES
add_textbox(s, 0.6, 1.2, 5.8, 0.4,
            "CMA-ES (Covariance Matrix Adaptation Evolution Strategy)",
            font_size=16, color=NAVY, bold=True)

cma_items = [
    "\u2022  Type: Evolutionary / black-box optimization",
    "\u2022  Few hyperparameters; handles continuous ranges",
    "\u2022  Achieves near-DQN performance (validates function design)",
    "",
    "Limitations:",
    "\u2022  Needs 24 episodes per parameter update (very slow)",
    "\u2022  Erratic exploration \u2014 unsafe during tuning",
    "\u2022  ~4,000 episodes to stabilize (11 years of simulated traffic)",
    "\u2022  Impractical for online / real-world deployment",
]
add_bullet_list(s, 0.8, 1.7, 5.5, 4.5, cma_items, font_size=14, color=DARK_GRAY, spacing=Pt(6))

# PPO
add_textbox(s, 7.0, 1.2, 5.8, 0.4,
            "PPO (Proximal Policy Optimization)",
            font_size=16, color=NAVY, bold=True)

ppo_items = [
    "\u2022  Type: Policy gradient method",
    "\u2022  Smooth, monotonic learning curves",
    "\u2022  Bounded gradient steps ensure safe exploration",
    "",
    "Limitations:",
    "\u2022  Gets stuck in local optima (over-regularized)",
    "\u2022  Suboptimal final convergence",
    "\u2022  Cannot beat actuated control under high demand",
    "\u2022  Insufficient for this domain despite safety benefits",
]
add_bullet_list(s, 7.2, 1.7, 5.5, 4.5, ppo_items, font_size=14, color=DARK_GRAY, spacing=Pt(6))

add_box(s, 0.6, 6.1, 12.1, 0.7,
        "Conclusion: Neither CMA-ES nor PPO are suitable for practical deployment. "
        "The DQN-based variants (especially DRHQ) are the most effective approach.",
        font_size=15, fill_color=BOX_BG, text_color=NAVY,
        border_color=NAVY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 8 — Experimental Setup
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Experimental Setup")

# Left
add_textbox(s, 0.6, 1.2, 5.8, 0.4, "Simulation Environment",
            font_size=17, color=NAVY, bold=True)
env_items = [
    "\u2022  Simulator: SUMO (Simulation of Urban Mobility)",
    "\u2022  Real traffic data from Utah DOT (2,092 intersections)",
    "\u2022  Test site: State St & E 4500 S, Murray, Utah",
    "\u2022  10 traffic phases, 11 non-conflicting phase pairs",
    "\u2022  >50,000 vehicles/day, peak rate: 95 cars/min",
    "\u2022  352 tunable parameters for regulatable policy",
]
add_bullet_list(s, 0.8, 1.7, 5.8, 3, env_items, font_size=14, color=DARK_GRAY, spacing=Pt(7))

# Right - demand profiles
add_textbox(s, 7.0, 1.2, 5.8, 0.4, "Traffic Demand Profiles",
            font_size=17, color=NAVY, bold=True)

add_box(s, 7.0, 1.7, 5.7, 0.55,
        "Low:      45,112 vehicles  |  1.04 v/sec  |  Wed May 1, 2019",
        font_size=13, fill_color=RGBColor(0xE8, 0xF5, 0xEB), text_color=DARK_GRAY,
        border_color=GREEN_ACC)
add_box(s, 7.0, 2.35, 5.7, 0.55,
        "Medium:  51,298 vehicles  |  1.19 v/sec  |  Mon May 6, 2019",
        font_size=13, fill_color=BOX_BG2, text_color=DARK_GRAY,
        border_color=RGBColor(0xD4, 0xA0, 0x60))
add_box(s, 7.0, 3.0, 5.7, 0.55,
        "High:     61,261 vehicles  |  1.42 v/sec  |  Fri Jun 21, 2019",
        font_size=13, fill_color=RGBColor(0xFC, 0xEB, 0xEB), text_color=DARK_GRAY,
        border_color=RED_ACC)

# DQN params
add_textbox(s, 7.0, 3.9, 5.8, 0.4, "DQN Hyperparameters",
            font_size=17, color=NAVY, bold=True)
dqn_items = [
    "\u2022  3 hidden layers, 64 units each",
    "\u2022  Replay buffer: 100,000 transitions",
    "\u2022  Minibatch size: 32  |  Optimizer: Adam",
    "\u2022  Epsilon-greedy: 0.05 \u2192 0 after 20 episodes",
    "\u2022  Discount factor: 0.8 (low/med), 0.9 (high)",
]
add_bullet_list(s, 7.2, 4.4, 5.5, 2.5, dqn_items, font_size=13, color=DARK_GRAY, spacing=Pt(5))

# Baseline
add_textbox(s, 0.6, 4.6, 5.8, 0.4, "Baseline: Actuated Signal Controller",
            font_size=17, color=NAVY, bold=True)
bl_items = [
    "\u2022  SUMO\u2019s actuated controller (used in real deployments)",
    "\u2022  Phases in fixed order (protected lefts \u2192 through traffic)",
    "\u2022  Maximum green time: 300 seconds",
    "\u2022  This is what real intersections use today",
]
add_bullet_list(s, 0.8, 5.1, 5.8, 2, bl_items, font_size=13, color=DARK_GRAY, spacing=Pt(5))

# ═══════════════════════════════════════════════
# SLIDE 9 — Main Results
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Main Results: Average Vehicle Delay Comparison")

add_box(s, 0.6, 1.2, 12.1, 0.55,
        "Average Vehicle Delay (seconds)  \u2014  Lower is Better",
        font_size=16, fill_color=NAVY, text_color=WHITE,
        border_color=NAVY, alignment=PP_ALIGN.CENTER)

# Table header
headers = "Method                         Low Demand        Medium Demand       High Demand"
add_textbox(s, 0.8, 1.95, 11.5, 0.35, headers,
            font_size=14, color=NAVY, bold=True, font_name="Consolas")

# Separator
add_bar(s, 0.8, 2.3, 11.5, 0.02, LIGHT_GRAY)

rows_data = [
    ("Actuated (Baseline)            ~60 sec            ~70 sec              ~95 sec", DARK_GRAY),
    ("CMA-ES                         ~50 sec            ~58 sec              ~85 sec", DARK_GRAY),
    ("PPO                            ~55 sec            ~65 sec              ~95 sec  (fails in high)", RED_ACC),
    ("DRQ                            Poor               Poor                 Poor", RED_ACC),
    ("DRSQ                           ~52 sec            ~62 sec              ~87 sec", DARK_GRAY),
    ("DRHQ  (Best Interpretable)     ~50 sec            ~60 sec              ~85 sec", GREEN_ACC),
    ("DQN   (Black-box Upper Bound)  ~48 sec            ~57 sec              ~82 sec", BLUE),
]

for i, (row, color) in enumerate(rows_data):
    y = 2.45 + i * 0.4
    add_textbox(s, 0.8, y, 11.5, 0.35, row,
                font_size=13, color=color, font_name="Consolas")
    if i < len(rows_data) - 1:
        add_bar(s, 0.8, y + 0.38, 11.5, 0.01, LIGHT_GRAY)

# Key finding
add_box(s, 0.6, 5.5, 12.1, 1.2,
        "Key Finding:  DRHQ achieves up to 19.4% reduced vehicle delay compared to deployed actuated controllers,\n"
        "while remaining fully interpretable.  The gap vs. the black-box DQN is only 1\u20133 seconds.",
        font_size=16, fill_color=RGBColor(0xE8, 0xF5, 0xEB), text_color=DARK_GRAY,
        border_color=GREEN_ACC, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 10 — Interpretation of Results
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Interpretation of Results")

add_textbox(s, 0.6, 1.2, 5.8, 0.4, "Why DRHQ Works Best",
            font_size=17, color=NAVY, bold=True)
drhq_items = [
    "\u2022  Only needs to match which action wins, not exact Q-values",
    "\u2022  Gives maximum flexibility to the polynomial function",
    "\u2022  Hardmax formulation directly targets policy equivalence",
    "\u2022  Converges within a single episode for low/medium demand",
    "\u2022  \"Less is more\": easier target = better practical fit",
]
add_bullet_list(s, 0.8, 1.7, 5.8, 3, drhq_items, font_size=15, color=DARK_GRAY, spacing=Pt(10))

add_textbox(s, 7.0, 1.2, 5.8, 0.4, "Why Other Methods Fall Short",
            font_size=17, color=NAVY, bold=True)
other_items = [
    "\u2022  DRQ: Cannot approximate the full Q-function surface",
    "\u2022  DRSQ: Softmax ranking harder to match than just the winner",
    "\u2022  PPO: Over-regularized, converges to local optima",
    "\u2022  CMA-ES: Requires thousands of episodes, unsafe exploration",
    "\u2022  Pure polynomial/Fourier: Cannot complete basic scenarios",
]
add_bullet_list(s, 7.2, 1.7, 5.8, 3, other_items, font_size=15, color=DARK_GRAY, spacing=Pt(10))

add_box(s, 0.6, 5.0, 12.1, 1.6,
        "Critical Insight\n\n"
        "The performance gap between DRHQ (~50\u201385 sec) and the full black-box DQN (~48\u201382 sec) is only 1\u20133 seconds.\n"
        "This marginal cost buys complete interpretability \u2014 every decision can be explained, audited,\n"
        "and manually adjusted by traffic engineers.  For safety-critical infrastructure, this is an excellent trade-off.",
        font_size=15, fill_color=BOX_BG, text_color=DARK_GRAY,
        border_color=NAVY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 11 — Limitations & Future Scope
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Limitations & Future Scope")

add_textbox(s, 0.6, 1.2, 5.8, 0.4, "Limitations of This Work",
            font_size=17, color=RED_ACC, bold=True)
lim_items = [
    "\u2022  Single intersection only \u2014 transferability is unknown",
    "\u2022  All results are from simulation (SUMO), no real-world test",
    "\u2022  CMA-ES is impractical despite showing ceiling performance",
    "\u2022  PPO fails under high demand scenarios",
    "\u2022  No warm-starting from existing controllers explored",
    "\u2022  352 parameters still requires significant compute",
    "\u2022  Limited to one geographic region (Utah traffic data only)",
]
add_bullet_list(s, 0.8, 1.7, 5.8, 4.5, lim_items, font_size=14, color=DARK_GRAY, spacing=Pt(8))

add_textbox(s, 7.0, 1.2, 5.8, 0.4, "Future Research Directions",
            font_size=17, color=GREEN_ACC, bold=True)
future_items = [
    "\u2022  Warm-starting from currently deployed controller behavior",
    "\u2022  Multi-intersection coordination & network optimization",
    "\u2022  Real-world deployment and field testing validation",
    "\u2022  Transferability studies across intersection types",
    "\u2022  Incorporating pedestrian and cyclist signal phases",
    "\u2022  Adaptive parameter count based on intersection complexity",
    "\u2022  Integration with connected & autonomous vehicle data",
]
add_bullet_list(s, 7.2, 1.7, 5.8, 4.5, future_items, font_size=14, color=DARK_GRAY, spacing=Pt(8))

# ═══════════════════════════════════════════════
# SLIDE 12 — Conclusion
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "Conclusion & Key Takeaways")

takeaways = [
    "1.   Interpretable polynomial policies CAN match DNN performance (gap \u2264 1\u20133 sec)",
    "",
    "2.   DRHQ (Hardmax variant) achieves 19.4% delay reduction vs. actuated controllers",
    "",
    "3.   \"Regulatable\" = monotonic input-output relationship = human-understandable decisions",
    "",
    "4.   Policy gradient methods (PPO) are unsuitable for traffic signal optimization",
    "",
    "5.   The paper bridges the gap between AI performance and real-world deployability",
    "",
    "6.   Interpretability is not a luxury \u2014 it is a requirement for safety-critical systems",
]
add_bullet_list(s, 1.0, 1.3, 11.3, 4.5, takeaways, font_size=17, color=DARK_GRAY, spacing=Pt(4))

add_box(s, 2.0, 5.8, 9.3, 0.9,
        "\"We can build AI systems that are both powerful and explainable.\n"
        "We do not always have to choose one over the other.\"",
        font_size=16, fill_color=BOX_BG, text_color=NAVY,
        border_color=NAVY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 13 — References
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_header(s, "References")

refs = [
    "[1]  Ault, J., Hanna, J. P., & Sharon, G. (2020). Learning an Interpretable Traffic Signal",
    "      Control Policy. arXiv preprint arXiv:1912.11023v2.",
    "",
    "[2]  Mnih, V. et al. (2015). Human-level control through deep reinforcement learning.",
    "      Nature, 518(7540), 529\u2013533.",
    "",
    "[3]  Schulman, J. et al. (2017). Proximal Policy Optimization Algorithms.",
    "      arXiv preprint arXiv:1707.06347.",
    "",
    "[4]  Hansen, N. (2006). The CMA Evolution Strategy: A Comparing Review.",
    "      Towards a New Evolutionary Computation, pp. 75\u2013102.",
    "",
    "[5]  Lopez, P.A. et al. (2018). Microscopic Traffic Simulation using SUMO.",
    "      IEEE ITSC 2018.",
]
add_bullet_list(s, 0.8, 1.3, 11.5, 5.5, refs, font_size=14, color=DARK_GRAY, spacing=Pt(3))

# ═══════════════════════════════════════════════
# SLIDE 14 — Thank You
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
add_bar(s, 0, 2.5, 13.333, 2.8, DARK_BLUE)

add_textbox(s, 1.5, 3.0, 10.3, 0.8, "Thank You",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_bar(s, 5.5, 3.95, 2.3, 0.03, RGBColor(0xBB, 0xCC, 0xDD))

add_textbox(s, 1.5, 4.2, 10.3, 0.6, "Questions & Discussion",
            font_size=22, color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 5.8, 10.3, 0.4,
            "Paper: arXiv 1912.11023v2  |  Authors: Ault, Hanna, Sharon (2020)",
            font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(s, 1.5, 6.3, 10.3, 0.4,
            "Presented by: Naaz Verma",
            font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_bar(s, 0, 7.2, 13.333, 0.03, BLUE)

# ═══════════════════════════════════════════════
output_path = r"c:\Users\naaz.verma\personal\python_projects\01_quiz_master\Traffic_Signal_Control_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
